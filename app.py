from __future__ import annotations
import base64
import os
import re
import time
import shutil
import threading
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET
from bisect import bisect_right
from concurrent.futures import ThreadPoolExecutor, as_completed
from io import BytesIO
from typing import Any, Dict, List, Tuple, Callable, Optional
from datetime import datetime
import tempfile
import queue
import logging
import sys
import unicodedata

import pandas as pd
import pdfplumber
import streamlit as st
import streamlit.components.v1 as components
from docx import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dataclasses import dataclass

# =========================
# Dependencias PDF / Word / PPT para procesamiento de documentos
# =========================
try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

try:
    import requests
except ImportError:
    requests = None

# ======================================================
# LOGGING
# ======================================================
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
logger = logging.getLogger(__name__)

# ======================================================
# CONFIG / CONSTANTES (Broken Link Checker)
# ======================================================
APP_TITLE = "UTP - GrammarScan"
APP_ICON = "📚"

MODULES = [
    "Home",
    "Report GrammarScan",
]

MAX_ZIP_BLOCK_MB = 200   # tamaño máximo aproximado por bloque de ZIP (MB)

# --------- Límite de seguridad Streamlit Cloud (solo afecta a la descarga masiva desde Excel) ----------
MAX_BULK_URLS_CLOUD = 700


def is_streamlit_cloud() -> bool:
    """Detecta si la app está ejecutándose en Streamlit Cloud.

    Se puede forzar de dos maneras:
    - Configurando st.secrets["is_streamlit_cloud"] = true/false.
    - Definiendo la variable de entorno IS_STREAMLIT_CLOUD=true/false.
    En local, por defecto devolverá False.
    """
    # 1) Intentar leer desde st.secrets (si existe)
    try:
        if "is_streamlit_cloud" in st.secrets:
            val = st.secrets["is_streamlit_cloud"]
        else:
            val = None
    except Exception:
        val = None

    # 2) Si no está en secrets, mirar variable de entorno
    if val is None:
        val = os.environ.get("IS_STREAMLIT_CLOUD", "")

    if isinstance(val, bool):
        return val

    return str(val).strip().lower() in ("1", "true", "yes", "y", "on")


IS_STREAMLIT_CLOUD = is_streamlit_cloud()


# Descarga Masiva
MAX_INTENTOS_DESCARGA = 7
CHUNK_SIZE = 1024 * 256
REQUEST_HEADERS = {
    "User-Agent": "Mozilla/5.0 (compatible; UTP-FileDownloader/1.0)",
    "Accept": "*/*",
    "Accept-Encoding": "identity",
}
DESC_EXT_PERMITIDAS = (".ppt", ".pptx", ".pdf", ".doc", ".docx")

# ======================================================
# LÓGICA GRAMMARSCAN (original)
# ======================================================
LINES_PER_TXT_PAGE = 50
PAGE_SEP = "\n\f\n"
ALLOWED_DOC_EXTS = {".pdf", ".docx", ".pptx", ".txt"}

WS_MULTI_RE = re.compile(r"[ \t]+")
NL_3PLUS_RE = re.compile(r"\n{3,}")

def normalize_ws(text: str) -> str:
    text = text.replace("\r", "\n")
    text = WS_MULTI_RE.sub(" ", text)
    text = NL_3PLUS_RE.sub("\n\n", text)
    return text.strip()

def find_java() -> bool:
    return shutil.which("java") is not None

def safe_str(x) -> str:
    try:
        return str(x)
    except Exception:
        return ""

ACCENTED_VOWELS = "áéíóúÁÉÍÓÚ"

def _has_accented_vowel(s: str) -> bool:
    return any(ch in ACCENTED_VOWELS for ch in s or "")

@dataclass
class ModismoPattern:
    modismo: str
    tipo: str
    patron: str
    sugerencia: str
    comentario: str
    regex: re.Pattern

@dataclass
class LogicalFileSource:
    display_name: str
    ext: str
    read_bytes: Callable[[], bytes]

def _normalize_regex_pattern(patron: str) -> str:
    if not isinstance(patron, str):
        patron = str(patron or "")
    patron = patron.replace("\\\\", "\\")
    return patron.strip()

def load_modismos_from_excel(path: str) -> List[ModismoPattern]:
    if not os.path.isfile(path):
        raise FileNotFoundError(f"No se encontró el archivo de modismos: {path}")

    df = pd.read_excel(path)
    required_cols = {"modismo", "tipo", "sugerencia"}
    missing = required_cols - set(df.columns)
    if missing:
        raise ValueError(f"Faltan columnas obligatorias en 'modismos_ar.xlsx': {missing}")

    patterns: List[ModismoPattern] = []
    for _, row in df.iterrows():
        modismo = safe_str(row.get("modismo", "")).strip()
        tipo = safe_str(row.get("tipo", "")).strip().lower()
        sugerencia = safe_str(row.get("sugerencia", "")).strip()
        comentario = safe_str(row.get("comentario", "")).strip()
        patron_cfg = safe_str(row.get("patron", "")).strip()

        if not modismo or not sugerencia:
            continue

        if tipo not in ("literal", "regex"):
            tipo = "literal"

        if tipo == "literal":
            patron = r"(?<!\w)" + re.escape(modismo) + r"(?!\w)"
        else:
            base = patron_cfg if patron_cfg else modismo
            patron = _normalize_regex_pattern(base)

        try:
            rx = re.compile(patron, flags=re.IGNORECASE | re.UNICODE)
        except re.error:
            continue

        patterns.append(
            ModismoPattern(
                modismo=modismo,
                tipo=tipo,
                patron=patron,
                sugerencia=sugerencia,
                comentario=comentario,
                regex=rx,
            )
        )

    return patterns

@st.cache_resource(show_spinner=False)
def get_modismos_patterns(modismos_path: str) -> List[ModismoPattern]:
    return load_modismos_from_excel(modismos_path)

# ======================================================
# DETECCIÓN ULTRA ROBUSTA DE BIBLIOGRAFÍAS / REFERENCIAS
# ======================================================

HEAD_RE = re.compile(
    r"^\s*("
    r"referencias?(\s+(bibliogr[aá]ficas?|consultadas?|citadas?))?|"
    r"bibliograf[ií]a(\s+(consultada|citada|utilizada))?|"
    r"bibliography|references?(\s+(cited|consulted))?|"
    r"works?\s+cited|obras?\s+citadas?|"
    r"literatura\s+citada|fuentes?\s+(bibliogr[aá]ficas?|consultadas?|de\s+consulta)|"
    r"webgraf[ií]a|webliograf[ií]a|netgraf[ií]a|"
    r"citas?\s+bibliogr[aá]ficas?|"
    r"literature\s+cited|cited\s+literature|"
    r"sources?(\s+consulted)?|consulted\s+sources?|"
    r"refer[eê]ncias(\s+bibliogr[aá]ficas?)?|"
    r"liste?\s+des?\s+r[ée]f[ée]rences|"
    r"literaturverzeichnis|quellenverzeichnis"
    r")\s*:?\s*$",
    re.IGNORECASE | re.UNICODE,
)

# --- URL base: dominios + esquemas conocidos ---
URL_BASE_RE = re.compile(
    r"("
    r"https?://[^\s]+"
    r"|www\.[^\s]+"
    r"|[A-Za-z0-9\-_.]+\.([A-Za-z]{2,})(/[^\s]*)?"
    r")",
    re.IGNORECASE,
)

# Partes válidas de un path de URL (para unir saltos de línea)
URL_PATH_CHUNK_RE = re.compile(r"[A-Za-z0-9\-._~:/?#\[\]@!$&'()*+,;=%]+")

def _merge_spans(spans: list[tuple[int, int]]) -> list[tuple[int, int]]:
    """Une spans solapados o adyacentes (p.ej. cuando unimos tramos de URL)."""
    if not spans:
        return []
    spans_sorted = sorted(spans, key=lambda s: s[0])
    merged: list[tuple[int, int]] = []
    cur_start, cur_end = spans_sorted[0]
    for start, end in spans_sorted[1:]:
        if start <= cur_end:  # solapa
            if end > cur_end:
                cur_end = end
        else:
            merged.append((cur_start, cur_end))
            cur_start, cur_end = start, end
    merged.append((cur_start, cur_end))
    return merged

def find_url_spans_in_text(text: str) -> list[tuple[int, int]]:
    """
    Devuelve spans (start, end) de URLs en el texto, uniendo casos como:
    https://reqtest.com/.../how-to-useinterviews-
    to-gather-requirements/

    Mantiene indices sobre el texto ORIGINAL.
    """
    spans: list[tuple[int, int]] = []
    if not text:
        return spans

    n = len(text)

    # 1) Encontrar matches iniciales (URLs base)
    for m in URL_BASE_RE.finditer(text):
        start, end = m.span()
        url_end = end

        # 2) Extender URL cuando hay:
        #    - salto de línea con '-' o '/'
        #    - seguido de fragmento válido de path
        while (
            url_end < n
            and text[url_end] == "\n"
            and url_end - 1 >= 0
            and text[url_end - 1] in "-/"
        ):
            pos = url_end + 1
            # Saltar espacios o tabs después del salto de línea
            while pos < n and text[pos] in (" ", "\t"):
                pos += 1
            cont_start = pos
            # Capturar siguiente token hasta el siguiente espacio/salto de línea
            while pos < n and not text[pos].isspace():
                pos += 1
            cont = text[cont_start:pos]
            if not cont:
                break
            # Verificamos que sea un fragmento válido de path de URL
            if not URL_PATH_CHUNK_RE.fullmatch(cont):
                break
            # Al menos debe tener algo tipo '/' o '.' para ser razonable
            if not any(ch in cont for ch in ("/", ".")):
                break

            url_end = pos  # extendemos el final del span

        spans.append((start, url_end))

    return _merge_spans(spans)

def mask_urls_preserve_length(
    text: str,
    mask_char: str = " ",
) -> tuple[str, list[tuple[int, int]]]:
    """
    Reemplaza TODOS los caracteres de los spans de URL por `mask_char`,
    pero conserva saltos de línea. Devuelve:
      - texto_enmascarado
      - lista de spans (start, end) de URLs en el texto original
    """
    spans = find_url_spans_in_text(text)
    if not spans:
        return text, []

    chars = list(text)
    for start, end in spans:
        for i in range(start, end):
            # Dejamos \n y \r tal cual para no romper estructura de páginas
            if chars[i] not in ("\n", "\r"):
                chars[i] = mask_char

    return "".join(chars), spans


DOI_URL_HINT_RE = re.compile(
    r"("
    r"doi:\s*10\.\d{4,9}/[^\s\)]+|"
    r"doi\.org/10\.\d{4,9}/[^\s\)]+|"
    r"urn:[^\s\)]+|"
    r"hdl\.handle\.net/[^\s\)]+|"
    r"pmid:\s*\d+|"
    r"arxiv:\s*[\d\.]+|"
    r"issn[\s:-]*\d{4}[\s-]?\d{3}[\dxX]|"
    r"isbn[\s:-]*(?:\d{9}[\dxX]|\d{13})"
    r")",
    re.IGNORECASE,
)

YEAR_RE = re.compile(r"\b(19|20)\d{2}[a-z]?\b")

JOURNAL_HINT_RE = re.compile(
    r"\b("
    r"vol(?:ume|umen)?\.?|"
    r"no\.?|n[úu]m(?:ero)?\.?|nº|"
    r"pp?\.?|p[aá]g(?:ina)?s?\.?|"
    r"ed(?:ición|ition)?\.?|"
    r"issue|"
    r"issn|isbn|"
    r"journal|revista|"
    r"proceedings|actas|"
    r"conference|congreso|"
    r"trans(?:action)?s?\.?"
    r")\b",
    re.IGNORECASE,
)

PUBLISHER_HINT_RE = re.compile(
    r"\b("
    r"press|editorial(?:es)?|ediciones?|"
    r"universidad(?:\s+(?:de|del|de\s+la))?\s+\w+|"
    r"university\s+(?:of\s+)?\w+|"
    r"pearson|mcgraw[- ]?hill|elsevier|springer|wiley|"
    r"cengage|prentice\s*hall|sage|oxford|cambridge|"
    r"harvard\s*(?:university\s*)?press|"
    r"mit\s*press|routledge|taylor\s*&?\s*francis|"
    r"blackwell|pergamon|academic\s*press|"
    r"john\s*wiley|norton|macmillan|"
    r"addison[- ]?wesley|thomson|"
    r"publi(?:shed|cado|cación)|"
    r"impreso\s+en|printed\s+in"
    r")\b",
    re.IGNORECASE,
)

ETAL_RE = re.compile(
    r"\b("
    r"et\s+al\.?|"
    r"et\s+alii|"
    r"y\s+(?:col(?:aboradores)?|cols?|otros)\.?|"
    r"and\s+(?:others|colleagues?)\.?|"
    r"e\s+(?:outros|cols?)\.?"
    r")\b",
    re.IGNORECASE,
)

BULLET_PREFIX_RE = re.compile(
    r"^\s*(?:"
    r"[\u2022\u2023\u25E6\u2043\u2219\u25CF\u25AA\u25AB\u25A0\u25A1]|"
    r"[-–—·•▪●◦‣]|"
    r"\[\d+\]|"
    r"\d+[\.\)]\s+"
    r")\s*"
)

def _strip_bullet(line: str) -> str:
    """
    Elimina viñetas y numeraciones de inicio de línea.
    """
    return BULLET_PREFIX_RE.sub("", line).strip()

# -------------------------------
# Patrones por estilo de cita
# -------------------------------

APA_AUTHOR_YEAR_RE = re.compile(
    r"("
    r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s+(?:[A-Z]\.?\s*){1,4}"
    r"(?:\s*(?:&|y|and|e)\s+[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s+(?:[A-Z]\.?\s*){1,4})*"
    r"(?:\s*(?:&|y|and|e)\s+[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s+(?:[A-Z]\.?\s*){1,4})*"
    r"\s*\(\d{4}[a-z]?\)"
    r")",
    re.UNICODE,
)

MLA_AUTHOR_RE = re.compile(
    r"^[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s+[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü\s]+\.",
    re.UNICODE,
)

IEEE_BRACKET_RE = re.compile(
    r"^\[\d+\]\s+[A-ZÁÉÍÓÚÑÜ]\.?\s+[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+",
    re.UNICODE,
)

VANCOUVER_NUM_RE = re.compile(
    r"^\d+\.\s+[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+\s+[A-Z](?:\s*[A-Z])?[\.,]",
    re.UNICODE,
)

AUTHOR_NAME_RE = re.compile(
    r"\b[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]{2,},\s*(?:[A-Z]\.?\s*){1,4}\b",
    re.UNICODE,
)

IN_TEXT_CITATION_RE = re.compile(
    r"\("
    r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+"
    r"(?:\s+(?:et\s+al\.?|y\s+(?:col\.?|otros)|and\s+others?))?"
    r"[\s,]+\d{4}[a-z]?"
    r"(?:\s*;\s*[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+(?:\s+(?:et\s+al\.?|y\s+col\.?))?\s*,?\s*\d{4}[a-z]?)*"
    r"\)",
    re.UNICODE,
)

BIB_RE_LIST = [
    APA_AUTHOR_YEAR_RE,
    MLA_AUTHOR_RE,
    IEEE_BRACKET_RE,
    VANCOUVER_NUM_RE,
    AUTHOR_NAME_RE,
    IN_TEXT_CITATION_RE,
]

def is_bibliography_heading(line: str) -> bool:
    """
    ¿La línea es un título de sección bibliográfica?
    """
    return bool(HEAD_RE.match(line.strip()))

def is_reference_line(line: str) -> bool:
    """
    Detecta si una línea aislada parece una referencia bibliográfica.
    """
    if not line or len(line.strip()) < 10:
        return False

    s = _strip_bullet(line)
    if not s:
        return False

    # 1) Patrones típicos de citas (APA, MLA, IEEE, etc.)
    for bib_re in BIB_RE_LIST:
        if bib_re.search(s):
            return True

    # 2) Heurísticas basadas en año + metadatos
    has_year = YEAR_RE.search(s)
    has_doi_urn = DOI_URL_HINT_RE.search(s)
    has_journal_meta = JOURNAL_HINT_RE.search(s)
    has_publisher = PUBLISHER_HINT_RE.search(s)
    has_url = URL_BASE_RE.search(s)
    has_etal = ETAL_RE.search(s)

    # Año + DOI/URN/URL
    if has_year and (has_doi_urn or has_url):
        return True

    # Año + metadatos de revista
    if has_year and has_journal_meta:
        return True

    # Año + editorial
    if has_year and has_publisher:
        return True

    # "et al." + año
    if has_etal and has_year:
        return True

    # Muchos autores separados por comas + año
    if has_year and s.count(",") >= 2:
        return True

    # Título entre comillas + año
    if has_year and re.search(r'["\'“”].+?["\'“”]\s*[\.,]', s):
        return True

    # Formatos IEEE/Vancouver con número o corchete al inicio
    if re.match(r"^(\[\d+\]|\d+[\.\)])\s+", s):
        if has_year or s.count(",") >= 1:
            return True

    # Múltiples patrones de "Apellido, I."
    author_pattern_count = len(
        re.findall(r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s*(?:[A-Z]\.?\s*)+", s)
    )
    if author_pattern_count >= 2:
        return True

    # Citas en texto estilo (Autor, 2020)
    if IN_TEXT_CITATION_RE.search(s):
        return True

    return False
def is_reference_fragment(text: str) -> bool:
    """
    Versión MEJORADA Y MÁS CONSERVADORA.
    Solo marca como referencia si hay MÚLTIPLES señales fuertes.
    """
    if not text or len(text.strip()) < 20:  # Aumentado de 10 a 20
        return False

    lines = text.splitlines()
    ref_line_count = 0

    for raw_line in lines:
        line = _strip_bullet(raw_line).strip()
        if not line or len(line) < 15:  # Aumentado de 10 a 15
            continue

        if is_reference_line(line):
            ref_line_count += 1

        if is_bibliography_heading(line):
            return True

    # CAMBIO CRÍTICO: Ahora requiere 3+ líneas (antes era 2+)
    if ref_line_count >= 3:
        return True

    # Análisis del texto plano completo
    flat = " ".join(_strip_bullet(s).strip() for s in lines if s.strip())

    # Contador de señales fuertes
    strong_signals = 0
    
    has_year = YEAR_RE.search(flat)
    has_doi = DOI_URL_HINT_RE.search(flat)
    has_journal = JOURNAL_HINT_RE.search(flat)
    has_publisher = PUBLISHER_HINT_RE.search(flat)
    has_etal = ETAL_RE.search(flat)
    has_url = URL_BASE_RE.search(flat)

    if has_year and has_doi:
        strong_signals += 2
    if has_year and has_url and 'doi' in flat.lower():
        strong_signals += 2
    if has_year and has_journal:
        strong_signals += 1
    if has_year and has_publisher and flat.count(",") >= 3:  # Más restrictivo
        strong_signals += 1
    if has_etal and has_year:
        strong_signals += 1

    # Múltiples autores + año
    author_count = len(
        re.findall(r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s*(?:[A-Z]\.?\s*)+", flat)
    )
    if author_count >= 3 and has_year:  # Aumentado de 2 a 3
        strong_signals += 1

    # CAMBIO CRÍTICO: Requiere 3+ señales fuertes (antes aceptaba con menos)
    return strong_signals >= 3

def has_author_pattern(text: str) -> bool:
    """
    Heurística adicional: ¿parece que el texto contiene nombres de autores típicos de bibliografía?
    """
    if not text:
        return False

    # Apellido, I.
    if re.search(r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s*(?:[A-Z]\.?\s*)+", text):
        return True

    # (Apellido, 2020) u otros patrones similares
    if IN_TEXT_CITATION_RE.search(text):
        return True

    # (2020) + apellido
    if re.search(r"\(\d{4}[a-z]?\)", text) and re.search(r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü]{2,}", text):
        return True

    # Apellido y Apellido + año
    apellidos_pattern = r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]{2,}"
    if re.search(rf"{apellidos_pattern}\s+(?:y|and|&|e)\s+{apellidos_pattern}", text):
        if YEAR_RE.search(text):
            return True

    return False
def detect_bibliography_pages(pages: List[Tuple[int, str]]) -> set:
    """
    Versión MEJORADA Y MÁS CONSERVADORA.
    Umbrales más altos para evitar falsos positivos.
    """
    skip: set[int] = set()
    in_bib_section = False
    consecutive_ref_pages = 0

    for idx, (num, txt) in enumerate(pages):
        lines_raw = [l.strip() for l in txt.splitlines() if l.strip()]
        if not lines_raw:
            continue

        lines = [_strip_bullet(l) for l in lines_raw if l.strip()]
        if not lines:
            continue

        # --- Encabezado en las primeras 5 líneas (antes 10) ---
        head_zone = lines[:5]
        has_heading = any(is_bibliography_heading(l) for l in head_zone)

        if has_heading:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages = 1
            logger.info(f"Página {num}: Encabezado de bibliografía detectado")
            continue

        # --- Métricas básicas de la página ---
        ref_count = sum(1 for l in lines if is_reference_line(l))
        url_count = sum(1 for l in lines if URL_BASE_RE.search(l))
        year_count = len(YEAR_RE.findall(" ".join(lines)))
        bullet_lines = sum(1 for l in lines_raw if BULLET_PREFIX_RE.match(l))
        etal_count = sum(1 for l in lines if ETAL_RE.search(l))
        doi_count = sum(1 for l in lines if DOI_URL_HINT_RE.search(l))
        journal_count = sum(1 for l in lines if JOURNAL_HINT_RE.search(l))
        publisher_count = sum(1 for l in lines if PUBLISHER_HINT_RE.search(l))
        author_pattern_count = sum(
            1
            for l in lines
            if re.search(r"[A-ZÁÉÍÓÚÑÜ][a-záéíóúñü'\-]+,\s*(?:[A-Z]\.?\s*)+", l)
        )

        n = max(1, len(lines))
        ref_ratio = ref_count / n

        # --- UMBRALES MÁS ALTOS PARA MARCAR COMO BIBLIOGRAFÍA ---
        
        # CAMBIO 1: Requiere 5+ referencias (antes 3+) o ratio >= 0.40 (antes 0.30)
        if ref_count >= 5 or ref_ratio >= 0.40:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por referencias ({ref_count})")
            continue

        # CAMBIO 2: Requiere 5+ URLs y 5+ años (antes 3+)
        if url_count >= 5 and year_count >= 5:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por URLs y años")
            continue

        # CAMBIO 3: Requiere 5+ viñetas y 4+ indicadores (antes 3+ y 2+)
        if bullet_lines >= 5 and (ref_count + url_count) >= 4:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por viñetas e indicadores")
            continue

        # CAMBIO 4: Requiere 4+ DOIs (antes 2+)
        if doi_count >= 4:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por DOIs")
            continue

        # CAMBIO 5: Requiere 4+ metadatos y 4+ años (antes 2+ y 2+)
        if journal_count >= 4 and year_count >= 4:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por metadatos de revista")
            continue

        # CAMBIO 6: Requiere 2+ publisher y 4+ años (antes 1+ y 2+)
        if publisher_count >= 2 and year_count >= 4:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por publisher")
            continue

        # CAMBIO 7: Requiere 5+ patrones de autor (antes 3+)
        if author_pattern_count >= 5:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por patrones de autor")
            continue

        # CAMBIO 8: Requiere 4+ etal y 4+ años (antes 2+ y 2+)
        if etal_count >= 4 and year_count >= 4:
            skip.add(num)
            in_bib_section = True
            consecutive_ref_pages += 1
            logger.info(f"Página {num}: Bibliografía por et al.")
            continue

        # --- Continuación de sección con criterios MÁS ESTRICTOS ---
        if in_bib_section:
            # CAMBIO 9: Criterios más estrictos para continuar sección
            if (
                (ref_count >= 2 and (url_count >= 2 or year_count >= 2))
                or (url_count >= 3 and year_count >= 2)
                or (ref_ratio >= 0.25)  # Aumentado de 0.20
                or (bullet_lines >= 3 and year_count >= 2)
                or (author_pattern_count >= 3)  # Aumentado de 2
            ):
                skip.add(num)
                consecutive_ref_pages += 1
                logger.info(f"Página {num}: Continuación de bibliografía")
            else:
                # Romper sección solo después de verificar siguiente página
                if consecutive_ref_pages >= 2:
                    if idx + 1 < len(pages):
                        next_num, next_txt = pages[idx + 1]
                        next_lines = [
                            _strip_bullet(l.strip())
                            for l in next_txt.splitlines()
                            if l.strip()
                        ]
                        next_ref_count = sum(
                            1 for l in next_lines if is_reference_line(l)
                        )
                        if next_ref_count < 2:  # Aumentado de 1
                            in_bib_section = False
                            consecutive_ref_pages = 0
                            logger.info(f"Página {num}: Fin de bibliografía")
                    else:
                        in_bib_section = False
                        consecutive_ref_pages = 0
                else:
                    consecutive_ref_pages = 0

    logger.info(f"Total páginas marcadas como bibliografía: {len(skip)}/{len(pages)}")
    return skip

EN_COMMON_WORDS = {
    "the", "and", "or", "but", "if", "then", "else", "when", "where", "who", "what",
    "which", "while", "for", "from", "to", "in", "on", "at", "by", "of", "with",
}

EN_TOKEN_RE = re.compile(r"[a-zA-Z]+")

def is_english_fragment(text: str) -> bool:
    if not text:
        return False

    cleaned = URL_BASE_RE.sub(" ", text)
    lower = cleaned.lower()

    if any(ch in lower for ch in "áéíóúüñ"):
        return False

    tokens = EN_TOKEN_RE.findall(lower)
    if len(tokens) < 4:
        return False

    en_hits = sum(1 for t in tokens if t in EN_COMMON_WORDS)
    if en_hits == 0:
        return False

    en_ratio = en_hits / max(1, len(tokens))
    ascii_tokens = [t for t in tokens if re.fullmatch(r"[a-z]+", t)]
    ascii_ratio = len(ascii_tokens) / max(1, len(tokens))

    if ascii_ratio >= 0.8 and en_ratio >= 0.25:
        return True
    if not any(ch in lower for ch in "áéíóúüñ") and en_ratio >= 0.35:
        return True

    return False

LATIN_KEYWORDS = {
    "ius", "honorarium", "civile", "sabinum", "edictum", "aedilium", "curulium",
    "officio", "pronoconsulis", "corpus", "delicti", "mens", "rea",
}

def is_latin_fragment(text: str) -> bool:
    if not text:
        return False

    cleaned = URL_BASE_RE.sub(" ", text)
    lower = cleaned.lower()
    tokens = EN_TOKEN_RE.findall(lower)
    if not tokens:
        return False

    hits = sum(1 for t in tokens if t in LATIN_KEYWORDS)

    if hits >= 1 and len(tokens) <= 4:
        return True
    if hits >= 2:
        return True
    if re.search(r"\bius\b", lower) and len(tokens) <= 6:
        return True

    return False

CODE_SYMBOLS = set("{}[]();<>:=+-*/%&|^#!@$\"'\\")
CODE_KEYWORDS = {
    "public ", "private ", "protected ", "class ", "interface ", "implements ",
    "extends ", "static ", "void ", "int ", "double ", "float ", "string ",
}

def is_code_fragment(text: str) -> bool:
    if not text:
        return False

    lines = [ln for ln in text.splitlines() if ln.strip()]
    if not lines:
        return False

    code_like_lines = 0
    total_nonempty = 0

    for ln in lines:
        s = ln.strip()
        if not s:
            continue
        total_nonempty += 1

        lower = s.lower()

        if any(kw in lower for kw in CODE_KEYWORDS):
            code_like_lines += 1
            continue

        if s.endswith((";", "{", "}")):
            code_like_lines += 1
            continue

        non_space = sum(1 for c in s if not c.isspace())
        if non_space >= 6:
            sym_count = sum(1 for c in s if c in CODE_SYMBOLS)
            if sym_count / non_space > 0.3:
                code_like_lines += 1
                continue

    if total_nonempty == 0:
        return False

    if code_like_lines >= 2 and code_like_lines / total_nonempty >= 0.5:
        return True

    all_non_space = sum(1 for c in text if not c.isspace())
    if all_non_space >= 10:
        all_sym = sum(1 for c in text if c in CODE_SYMBOLS)
        if all_sym / all_non_space > 0.35:
            return True

    return False

def read_pdf_pages(bio: BytesIO) -> List[Tuple[int, str]]:
    pages: List[Tuple[int, str]] = []
    with pdfplumber.open(bio) as pdf:
        for i, p in enumerate(pdf.pages, start=1):
            t = p.extract_text() or ""
            t = normalize_ws(t)
            if t:
                pages.append((i, t))
    return pages

def _iter_block_items(doc: Document):
    body = doc.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)

def _paragraph_has_page_break(para: Paragraph) -> bool:
    try:
        if para.paragraph_format.page_break_before:
            return True
    except Exception:
        pass
    try:
        for run in para.runs:
            if run._element.xpath('.//w:br[@w:type="page"]'):
                return True
    except Exception:
        pass
    return False

def read_docx_pages(bio: BytesIO) -> List[Tuple[int, str]]:
    doc = Document(bio)
    pages: List[Tuple[int, str]] = []
    buff: List[str] = []
    page_no = 1

    def flush():
        nonlocal buff, page_no, pages
        txt = normalize_ws("\n".join(buff))
        if txt:
            pages.append((page_no, txt))
        buff = []

    for b in _iter_block_items(doc):
        if isinstance(b, Paragraph):
            t = b.text.strip()
            if t:
                buff.append(t)
            if _paragraph_has_page_break(b):
                flush()
                page_no += 1
        else:
            for row in b.rows:
                row_text = " | ".join(
                    (c.text or "").strip() for c in row.cells if (c.text or "").strip()
                )
                if row_text:
                    buff.append(row_text)

    flush()

    if not pages:
        flat = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if flat.strip():
            pages = [(1, normalize_ws(flat))]

    return pages

def _clean_reference_lines_block(text: str) -> str:
    kept = []
    for ln in (text or "").splitlines():
        s = normalize_ws(ln)
        if not s:
            continue
        if (
            is_bibliography_heading(s)
            or is_reference_line(s)
            or (PUBLISHER_HINT_RE.search(s) and YEAR_RE.search(s))
            or is_code_fragment(s)
        ):
            continue
        kept.append(s)
    return "\n".join(kept).strip()

def _read_pptx_via_zip(bio: BytesIO) -> List[Tuple[int, str]]:
    slides: List[Tuple[int, str]] = []
    bio.seek(0)
    with zipfile.ZipFile(bio) as z:
        names = [
            n for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        ]

        def slide_no(name: str) -> int:
            m = re.search(r"slide(\d+)\.xml$", name)
            return int(m.group(1)) if m else 0

        names.sort(key=slide_no)

        A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

        for idx, name in enumerate(names, start=1):
            try:
                xml_bytes = z.read(name)
                root = ET.fromstring(xml_bytes)
                texts = [t.text for t in root.iter(f"{A_NS}t") if t.text]
                raw = "\n".join(texts)
                cleaned = _clean_reference_lines_block(normalize_ws(raw))
                if cleaned:
                    slides.append((idx, cleaned))
            except Exception:
                continue
    return slides

def _iter_shape_texts(shape) -> List[str]:
    textos: List[str] = []

    try:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            for para in shape.text_frame.paragraphs:
                frags = [run.text for run in para.runs if run.text]
                if frags:
                    textos.append("".join(frags))
    except Exception:
        pass

    try:
        if getattr(shape, "has_table", False) or getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.TABLE:
            tbl = shape.table
            for row in tbl.rows:
                celdas = []
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        celdas.append(t)
                if celdas:
                    textos.append(" | ".join(celdas))
    except Exception:
        pass

    try:
        subshapes = getattr(shape, "shapes", None)
        if subshapes is not None:
            for sub in subshapes:
                textos.extend(_iter_shape_texts(sub))
    except Exception:
        pass

    return textos

def read_pptx_slides(bio: BytesIO) -> List[Tuple[int, str]]:
    try:
        prs = Presentation(bio)
        slides: List[Tuple[int, str]] = []
        slide_h = float(prs.slide_height) if hasattr(prs, "slide_height") else None

        for s_idx, slide in enumerate(prs.slides, start=1):
            chunk: List[str] = []

            for sh in slide.shapes:
                for raw in _iter_shape_texts(sh):
                    if not raw:
                        continue

                    raw_norm = normalize_ws(raw)

                    try:
                        if slide_h and float(getattr(sh, "top", 0)) >= 0.75 * slide_h:
                            if (
                                is_reference_fragment(raw_norm)
                                or URL_BASE_RE.search(raw_norm)
                                or (
                                    PUBLISHER_HINT_RE.search(raw_norm)
                                    and YEAR_RE.search(raw_norm)
                                )
                            ):
                                continue
                    except Exception:
                        pass

                    cleaned = _clean_reference_lines_block(raw_norm)
                    if cleaned:
                        chunk.append(cleaned)

            txt = normalize_ws("\n".join(chunk))
            if txt:
                slides.append((s_idx, txt))
        return slides

    except Exception:
        try:
            return _read_pptx_via_zip(bio)
        except Exception as e2:
            raise e2

def read_txt_pages(bio: BytesIO) -> List[Tuple[int, str]]:
    try:
        raw = bio.read()
        for enc in ("utf-8-sig", "utf-8", "latin-1"):
            try:
                s = raw.decode(enc)
                break
            except Exception:
                continue
        else:
            s = raw.decode("utf-8", errors="ignore")
    except Exception:
        return []

    s = s.replace("\r\n", "\n").replace("\r", "\n")
    lines = s.split("\n")

    pages: List[Tuple[int, str]] = []
    for i in range(0, len(lines), LINES_PER_TXT_PAGE):
        page_num = (i // LINES_PER_TXT_PAGE) + 1
        chunk = normalize_ws("\n".join(lines[i:i + LINES_PER_TXT_PAGE]))
        if chunk:
            pages.append((page_num, chunk))

    if not pages and s.strip():
        pages = [(1, normalize_ws(s))]

    return pages

def extract_pages(file_bytes: bytes, file_name: str) -> Tuple[List[Tuple[int, str]], str]:
    ext = os.path.splitext(file_name)[1].lower()
    bio = BytesIO(file_bytes)

    if ext == ".pdf":
        return read_pdf_pages(bio), "Página"
    if ext == ".docx":
        return read_docx_pages(bio), "Página"
    if ext == ".pptx":
        return read_pptx_slides(bio), "Diapositiva"
    if ext == ".txt":
        return read_txt_pages(bio), "Página"

    return [], "Página"
def build_global_text(
    pages: List[Tuple[int, str]]
) -> Tuple[str, List[int], List[Tuple[int, int, int]]]:
    """
    Construye:
    - texto_global: concatenación de todas las páginas con separador PAGE_SEP
    - starts: lista con el offset inicial (en texto_global) de cada página
    - bounds: lista de tuplas (start, end, page_number) sincronizada con starts

    IMPORTANTE:
    - starts[i] y bounds[i] SIEMPRE corresponden a la MISMA página.
    - El separador PAGE_SEP solo se añade entre páginas, nunca antes de la primera
      ni después de la última.
    """
    parts: List[str] = []
    starts: List[int] = []
    bounds: List[Tuple[int, int, int]] = []

    cur = 0
    for idx, (num, txt) in enumerate(pages):
        # Offset donde empieza esta página en el texto global
        starts.append(cur)
        start = cur

        # Texto de la página
        parts.append(txt)
        cur += len(txt)

        end = cur
        bounds.append((start, end, num))

        # Añadir separador entre páginas (NO después de la última)
        if idx < len(pages) - 1:
            parts.append(PAGE_SEP)
            cur += len(PAGE_SEP)

    return "".join(parts), starts, bounds

def chunk_by_pages_with_offsets(
    pages: List[Tuple[int, str]],
    starts: List[int],
    bounds: List[Tuple[int, int, int]],
    max_chars: int,
) -> List[Tuple[int, int, int, int]]:
    """
    Devuelve una lista de chunks, cada uno como:
        (chunk_start_offset, chunk_end_offset, first_page_idx, last_page_idx_exclusive)

    - chunk_start_offset / chunk_end_offset: offsets en texto_global (coherentes con `starts` y `bounds`)
    - first_page_idx / last_page_idx_exclusive: índices de página en `pages` (0-based, tipo range(i, j))
    """
    chunks: List[Tuple[int, int, int, int]] = []
    n = len(pages)
    if n == 0:
        return chunks

    i = 0
    while i < n:
        # Primera página del chunk
        first_idx = i
        chunk_start = starts[first_idx]

        j = i
        last_offset = chunk_start
        while j < n:
            page_start, page_end, _ = bounds[j]
            page_len = page_end - page_start

            # Si no es la primera página, añadimos el separador PAGE_SEP
            extra = page_len
            if j > first_idx:
                extra += len(PAGE_SEP)

            # Longitud actual del chunk (desde chunk_start hasta last_offset)
            current_len = last_offset - chunk_start
            if current_len + extra > max_chars and j > first_idx:
                break

            # Aceptamos esta página en el chunk
            last_offset = page_end
            j += 1

        # last_offset es el final del chunk (en texto_global)
        chunks.append((chunk_start, last_offset, first_idx, j))
        i = j

    return chunks

from bisect import bisect_right  # ya lo tienes importado arriba

def page_for_offset(
    starts: List[int],
    bounds: List[Tuple[int, int, int]],
    offset: int
) -> int:
    """
    Dado un offset en el texto global, devuelve el número de página/diapositiva.

    - starts[i]  = offset donde empieza el texto de la página i (en texto_global)
    - bounds[i]  = (start_text, end_text, page_num) para esa misma página

    Usamos bisect_right sobre starts para localizar la página en O(log n).
    Si el offset cae en el separador PAGE_SEP, se asigna a la página anterior.
    """
    if not starts or not bounds:
        return 1

    # Normalizar offset extremo inferior
    if offset <= starts[0]:
        idx = 0
    else:
        # bisect_right devuelve la posición de inserción, restamos 1 para obtener
        # el índice de la página cuyo inicio es el último <= offset
        idx = bisect_right(starts, offset) - 1
        if idx < 0:
            idx = 0
        elif idx >= len(starts):
            idx = len(starts) - 1

    return bounds[idx][2]


@st.cache_resource(show_spinner=False)
def get_language_tool(lang_code: str):
    """
    Versión CORREGIDA con todas las reglas activas por defecto.
    
    IMPORTANTE: No intentamos habilitar reglas específicas porque:
    1. Los nombres de las reglas pueden variar entre versiones
    2. Muchas reglas ya están habilitadas por defecto
    3. Es más seguro DESACTIVAR solo las problemáticas
    """
    if not find_java():
        raise RuntimeError("Java no detectado. Activa tu JRE/JDK para usar LanguageTool local.")
    
    import language_tool_python as lt
    
    # Crear instancia básica - LanguageTool ya viene con TODAS las reglas activas
    tool = lt.LanguageTool(lang_code)
    # Ver cuántas reglas están activas
    logger.info(f"Reglas activas: {len(tool.enabled_rules)}")
    for rule in list(tool.enabled_rules)[:10]:
        logger.info(f"  - {rule.ruleId}")
    
    # Solo desactivamos reglas que causan falsos positivos
    rules_to_disable = [
        # Reglas de formato que generan ruido
        # 'WHITESPACE_RULE',           # Espacios en blanco
        # 'DOUBLE_PUNCTUATION',        # Puntuación doble
        
        # NO desactivamos UPPERCASE_SENTENCE_START para que detecte mayúsculas
    ]
    
    for rule in rules_to_disable:
        try:
            tool.disable_rule(rule)
        except Exception:
            pass
    
    logger.info(f"LanguageTool configurado: {len(tool.enabled_rules)} reglas activas")
    
    return tool


LT_LOCK = threading.Lock()

def analyze_text(tool, text: str, retries: int = 2, sleep: float = 0.8) -> list:
    if not text.strip():
        return []
    last_err = None
    for _ in range(retries + 1):
        try:
            with LT_LOCK:
                return tool.check(text)
        except Exception as e:
            last_err = e
            msg = str(e)
            if "stdin" in msg or "Connection refused" in msg or "WinError" in msg:
                try:
                    tool = get_language_tool(st.session_state.get("_lang_code", "es"))
                except Exception:
                    pass
            time.sleep(sleep)
    raise RuntimeError(f"Fallo LanguageTool local tras reintentos: {last_err}")

def detect_modismos_in_pages(
    file_name: str,
    pages: List[Tuple[int, str]],
    unit_label: str,
    patterns: List[ModismoPattern],
    skip_pages: set | None = None,
) -> pd.DataFrame:
    if not patterns:
        return pd.DataFrame([])

    skip_pages = skip_pages or set()
    rows: List[Dict[str, Any]] = []

    for page_no, text in pages:
        if page_no in skip_pages:
            continue
        if not text.strip():
            continue

        lower = text.lower()

        for pat in patterns:
            for m in pat.regex.finditer(lower):
                start, end = m.span()
                ctx_start = max(0, start - 60)
                ctx_end = min(len(text), end + 60)
                contexto = text[ctx_start:ctx_end]

                if is_reference_fragment(contexto):
                    continue

                match_text = text[start:end]

                if _has_accented_vowel(pat.modismo) and not _has_accented_vowel(match_text):
                    continue

                mensaje = (
                    f"Uso de modismo argentino «{match_text}». "
                    f"Sugerencia: «{pat.sugerencia}»."
                )

                rows.append({
                    "Archivo": file_name,
                    "Página/Diapositiva": page_no,
                    "BloqueTipo": unit_label,
                    "Mensaje": mensaje,
                    "Sugerencias": pat.sugerencia,
                    "Oración": contexto,
                    "Contexto": contexto,
                    "Regla": "MODISMO_AR",
                    "Categoría": f"UTP_CUSTOM: Modismos argentinos ({pat.modismo})",
                })

    if not rows:
        return pd.DataFrame([])

    df_mod = pd.DataFrame.from_records(rows).drop_duplicates()
    return df_mod

def analyze_file(
    file_name: str,
    file_bytes: bytes,
    lang_code: str,
    max_chars_call: int,
    workers: int,
    excluir_bibliografia: bool = True,
    modismos_patterns: List[ModismoPattern] | None = None,
    analizar_modismos: bool = False,
) -> pd.DataFrame:
    # 1) Extraer páginas/diapositivas
    pages, unit_label = extract_pages(file_bytes, file_name)
    if not pages:
        return pd.DataFrame([])

    logger.info("=" * 80)
    logger.info(f"Iniciando análisis de: {file_name}")
    logger.info(f"Idioma: {lang_code}")
    logger.info(f"Excluir bibliografía: {excluir_bibliografia}")
    logger.info(f"Analizar modismos: {analizar_modismos}")
    logger.info("=" * 80)

    # 2) Ajustar dinámicamente max_chars_call según el tamaño del documento
    total_chars = sum(len(txt) for _, txt in pages)
    target_blocks = 14
    approx_chars_per_block = max(5000, min(20000, total_chars // max(1, target_blocks)))
    if approx_chars_per_block < max_chars_call:
        max_chars_call = approx_chars_per_block

    ext = os.path.splitext(file_name)[1].lower()

    # 3) Detectar páginas de bibliografía (para filtros posteriores)
    skip_pages = detect_bibliography_pages(pages) if excluir_bibliografia else set()
    if excluir_bibliografia and ext in (".txt", ".docx", ".pptx"):
        # Para estos formatos, preferimos NO saltarnos páginas completas,
        # solo filtrar por contexto bibliográfico.
        skip_pages = set()

    # 4) Construir texto global y rangos de páginas (SINCRONIZADOS)
    st.session_state["_lang_code"] = lang_code
    texto_global, starts, bounds = build_global_text(pages)

    # Chunks basados en offsets globales para mantener coherencia con `starts` y `bounds`
    ranges = chunk_by_pages_with_offsets(
        pages=pages,
        starts=starts,
        bounds=bounds,
        max_chars=max_chars_call,
    )

    # 5) Instanciar LanguageTool
    tool = get_language_tool(lang_code)


    rows: List[Dict[str, Any]] = []
    lock = threading.Lock()

    total_ranges = len(ranges)
    done_ranges = 0

    max_effective_workers = max(1, min(workers, 3))
    progress_chunks_ph = st.empty()

    def _update_chunk_progress():
        if total_ranges <= 0:
            return
        pct = done_ranges / total_ranges
        detail = f"{done_ranges}/{total_ranges} bloques analizados"
        try:
            render_task_progress(
                progress_chunks_ph,
                "Analizando bloques de texto (LanguageTool)",
                pct,
                detail,
            )
        except Exception:
            pass

    def _extract_sentence(chunk_text: str, offset: int) -> str:
        """
        Extrae una 'oración' aproximada alrededor del error dentro del chunk.
        Busca saltos de línea o separadores de oración.
        """
        # Buscar inicio de oración en el chunk
        start = max(0, offset - 200)
        end = min(len(chunk_text), offset + 200)

        ventana = chunk_text[start:end]

        # Cortar por saltos de línea para hacerlo más legible
        lineas = ventana.split("\n")
        oracion = " ".join(l.strip() for l in lineas if l.strip())

        if len(oracion) > 250:
            oracion = oracion[:250].rstrip() + "..."

        return oracion
    def worker(rng: Tuple[int, int, int, int]):
        """
        Procesa un rango de páginas y devuelve las incidencias detectadas
        en ese chunk.

        rng = (chunk_start_offset, chunk_end_offset, first_page_idx, last_page_idx_exclusive)
        """
        nonlocal done_ranges

        chunk_start_offset, chunk_end_offset, first_idx, last_idx_exclusive = rng

        # Texto ORIGINAL del chunk (para contexto / oraciones)
        chunk_text_original = texto_global[chunk_start_offset:chunk_end_offset]

        if not chunk_text_original.strip():
            with lock:
                done_ranges += 1
                _update_chunk_progress()
            return []

        # 🔐 NUEVO: enmascarar URLs para que NO pasen a LanguageTool
        # Conserva longitud y saltos de línea ⇒ offsets siguen siendo válidos
        chunk_text_masked, url_spans = mask_urls_preserve_length(chunk_text_original)

        # Analizar con LanguageTool usando el texto enmascarado
        matches = analyze_text(tool, chunk_text_masked)

        local_rows: List[Dict[str, Any]] = []

        for m in matches:
            offset_in_chunk = m.offset

            # Offset global = inicio global del chunk + offset local
            global_offset = chunk_start_offset + offset_in_chunk

            # Determinar la página real con el offset global
            page_num = page_for_offset(starts, bounds, global_offset)

            # Extraer contexto desde el TEXTO ORIGINAL (para que el usuario vea el texto real)
            ctx_start = max(0, offset_in_chunk - 120)
            ctx_end = min(len(chunk_text_original), offset_in_chunk + m.errorLength + 120)
            contexto = chunk_text_original[ctx_start:ctx_end]

            # Oración más legible dentro del chunk (usamos también el texto original)
            oracion = _extract_sentence(chunk_text_original, offset_in_chunk)

            # Evitar páginas marcadas como bibliografía (solo si se usa skip_pages)
            if excluir_bibliografia and page_num in skip_pages:
                continue

            mensaje = m.message or "Error gramatical o de estilo"
            sugerencias = ", ".join(m.replacements[:3]) if m.replacements else ""
            regla = m.ruleId or "UNKNOWN"
            categoria = getattr(m, "category", None) or "General"

            local_rows.append({
                "Archivo": file_name,
                "Página/Diapositiva": page_num,
                "BloqueTipo": unit_label,
                "Mensaje": mensaje,
                "Sugerencias": sugerencias,
                "Oración": oracion,
                "Contexto": contexto,
                "Regla": regla,
                "Categoría": categoria,
            })

        with lock:
            done_ranges += 1
            _update_chunk_progress()

        return local_rows



    # 6) Ejecutar workers (paralelismo a nivel de chunks)
    with ThreadPoolExecutor(max_workers=max_effective_workers) as ex:
        futures = [ex.submit(worker, r) for r in ranges]
        for fut in as_completed(futures):
            part = fut.result()
            with lock:
                rows.extend(part)

    # 7) Construir DataFrame principal de LanguageTool
    if not rows:
        df_lt = pd.DataFrame([])
    else:
        df_lt = pd.DataFrame.from_records(rows)

        # 🔍 Filtro extra: eliminar incidencias cuya "Oración" es prácticamente solo un enlace
        def _looks_like_url_only(text: str) -> bool:
            if not text:
                return False
            t = text.strip()
            # Quitar comillas/puntos alrededor
            t = t.strip("()[]{}<>\"'.,;: ")
            if len(t) < 8:
                return False
            m = URL_BASE_RE.fullmatch(t)
            if m:
                return True
            # También casos donde el enlace es >70% del texto
            urls = URL_BASE_RE.findall(t)
            if not urls:
                return False
            total_len = len(t)
            urls_len = sum(len(u[0]) if isinstance(u, tuple) else len(u) for u in urls)
            return urls_len / max(1, total_len) > 0.7

        df_lt = df_lt.loc[~df_lt["Oración"].fillna("").apply(_looks_like_url_only)].copy()


    # 8) Filtros avanzados: bibliografía, código, inglés
    if not df_lt.empty:
        initial_count = len(df_lt)
        logger.info(f"Errores detectados inicialmente: {initial_count}")

        # --------------------------
        # FILTRO 1: bibliografía
        # --------------------------
        if excluir_bibliografia:
            before_filtro1 = len(df_lt)

            def is_clearly_bibliographic(text: str) -> bool:
                if not text or len(text) < 30:
                    return False

                # IEEE "[1] A. Author..."
                if re.search(r'^\[?\d+\]?\s+[A-Z]\.\s+[A-ZÁÉÍÓÚÑa-záéíóúñ]+.*\d{4}', text):
                    return True

                # Vancouver "1. Autor J..."
                if re.search(r'^\d+\.\s+[A-ZÁÉÍÓÚÑa-záéíóúñ]+\s+[A-Z]\b.*\d{4}', text):
                    return True

                # [1] ... and 2020
                if re.search(r'\[\d+\].*\b(and|y|&)\b.*\d{4}', text, re.I):
                    return True

                # Palabras clave de publicación
                pub_keywords = [
                    r'\bProc\.', r'\bConf\.', r'\bInt\.', r'\bTrans\.', r'\bJ\.',
                    r'\bvol\.', r'\bno\.', r'\bpp\.', r'\bEd\.',
                ]
                kw_count = sum(1 for kw in pub_keywords if re.search(kw, text, re.I))
                has_year = bool(re.search(r'\b(19|20)\d{2}\b', text))

                if kw_count >= 2 and has_year:
                    return True

                # DOI/URL académica
                if re.search(r'(doi:|https?://doi\.org/)', text, re.I):
                    return True

                # Múltiples autores "Apellido, I."
                author_pattern = r'[A-ZÁÉÍÓÚÑa-záéíóúñ]+,\s+[A-Z]\.'
                author_count = len(re.findall(author_pattern, text))
                if author_count >= 2 and has_year:
                    return True

                signals = 0
                if re.search(r'\(\d{4}[a-z]?\)', text):
                    signals += 2
                if re.search(r'\bet\s+al\.', text, re.I):
                    signals += 2
                if re.search(r'["“].+["”]', text):
                    signals += 1
                if re.search(r'\bpp?\.\s*\d+[-–]\d+', text, re.I):
                    signals += 2

                return signals >= 3

            mask_bib = (
                df_lt["Oración"].fillna("").apply(is_clearly_bibliographic) |
                df_lt["Contexto"].fillna("").apply(is_clearly_bibliographic)
            )
            df_lt = df_lt.loc[~mask_bib].copy()
            filtered = before_filtro1 - len(df_lt)
            logger.info(f"Filtro bibliografía: {filtered} errores eliminados")

        # --------------------------
        # FILTRO 2: código
        # --------------------------
        before_filtro2 = len(df_lt)

        def is_obvious_code(text: str) -> bool:
            if not text or len(text) < 15:
                return False
            if re.search(r'\b(function|class|def|var|const|import|return)\b', text):
                return True
            symbols = sum(1 for c in text if c in '{}[]();=<>:/*')
            if len(text) > 10 and symbols / len(text) > 0.35:
                return True
            if re.search(r'\)\s*;?\s*$', text.strip()):
                return True
            return False

        mask_code = (
            df_lt["Oración"].fillna("").apply(is_obvious_code) |
            df_lt["Contexto"].fillna("").apply(is_obvious_code)
        )
        df_lt = df_lt.loc[~mask_code].copy()
        filtered = before_filtro2 - len(df_lt)
        logger.info(f"Filtro código: {filtered} errores eliminados")

        # --------------------------
        # FILTRO 3: inglés puro (si idioma base es español)
        # --------------------------
        if lang_code.startswith("es"):
            before_filtro3 = len(df_lt)

            def is_pure_english_text(text: str) -> bool:
                if not text or len(text) < 30:
                    return False
                if re.search(r'[áéíóúüñ]', text, re.I):
                    return False
                english_words = {
                    'the', 'and', 'or', 'in', 'on', 'at', 'to', 'for', 'with',
                    'by', 'from', 'this', 'that', 'these', 'those', 'is', 'are'
                }
                words = re.findall(r"[a-zA-Z]+", text.lower())
                if len(words) < 5:
                    return False
                english_count = sum(1 for w in words if w in english_words)
                return english_count / len(words) > 0.5

            mask_eng = (
                df_lt["Oración"].fillna("").apply(is_pure_english_text) |
                df_lt["Contexto"].fillna("").apply(is_pure_english_text)
            )
            df_lt = df_lt.loc[~mask_eng].copy()
            filtered = before_filtro3 - len(df_lt)
            logger.info(f"Filtro inglés: {filtered} errores eliminados")

        # --------------------------
        # FILTRO 4: páginas completas marcadas como bibliografía
        # --------------------------
        if excluir_bibliografia and "Página/Diapositiva" in df_lt.columns and skip_pages:
            before_filtro4 = len(df_lt)
            mask_bib_page = df_lt["Página/Diapositiva"].isin(skip_pages)
            df_lt = df_lt.loc[~mask_bib_page].copy()
            filtered = before_filtro4 - len(df_lt)
            logger.info(f"Filtro páginas bibliográficas: {filtered} errores eliminados")

        final_count = len(df_lt)
        logger.info(f"Errores finales: {final_count} (de {initial_count} iniciales)")
        if initial_count > 0:
            logger.info(f"Tasa de retención: {final_count / initial_count * 100:.1f}%")

    # 9) Modismos argentinos (si corresponde)
    if analizar_modismos and lang_code.startswith("es") and modismos_patterns:
        if excluir_bibliografia and ext in (".txt", ".docx", ".pptx"):
            skip_for_modismos = None
        else:
            skip_for_modismos = skip_pages

        df_mod = detect_modismos_in_pages(
            file_name=file_name,
            pages=pages,
            unit_label=unit_label,
            patterns=modismos_patterns,
            skip_pages=skip_for_modismos,
        )
    else:
        df_mod = pd.DataFrame([])

    # 10) Combinar resultados LT + modismos
    if df_lt is None or df_lt.empty:
        final_df = df_mod
    elif df_mod is None or df_mod.empty:
        final_df = df_lt
    else:
        final_df = pd.concat([df_lt, df_mod], ignore_index=True)

    if final_df is None or final_df.empty:
        return pd.DataFrame([])

    final_df.sort_values(["Archivo", "Página/Diapositiva"], inplace=True)
    final_df.drop_duplicates(
        subset=["Archivo", "Página/Diapositiva", "Mensaje", "Oración", "Contexto", "Regla"],
        keep="first",
        inplace=True,
    )
    final_df.reset_index(drop=True, inplace=True)
    return final_df


# Patrones para limpiar caracteres no válidos en XML/Excel
INVALID_XML_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F]")

def _sanitize_excel_df(df: Optional[pd.DataFrame]) -> Optional[pd.DataFrame]:
    """
    Elimina de un DataFrame los caracteres de control no permitidos por
    el formato XLSX (XML 1.0) para evitar que Excel muestre errores
    de reparación al abrir el archivo.
    """
    if df is None or df.empty:
        return df

    df = df.copy()

    # Limpiar nombres de columnas
    df.columns = [INVALID_XML_RE.sub("", str(c)) for c in df.columns]

    # Limpiar el contenido de las celdas de tipo string
    def _clean_value(v: Any) -> Any:
        if isinstance(v, str):
            return INVALID_XML_RE.sub("", v)
        return v

    return df.applymap(_clean_value)


# 🔹 NUEVO: helper para detectar celdas vacías o "0"
def _is_empty_or_zero_cell(v: Any) -> bool:
    """
    Devuelve True si el valor debe considerarse vacío para efectos del reporte:
    - None, NaN
    - cadena vacía
    - cadena "0"
    - número 0
    """
    if v is None:
        return True

    # Números
    if isinstance(v, (int, float)):
        if pd.isna(v):
            return True
        return v == 0

    # Cadenas u otros tipos convertidos a string
    s = str(v).strip()
    if s == "" or s == "0":
        return True
    if s.lower() in ("nan", "none"):
        return True

    return False


def _filter_resultados_empty_suggest_or_sentence(
    df: Optional[pd.DataFrame],
) -> Optional[pd.DataFrame]:
    """
    Elimina filas donde:
    - Sugerencias esté vacía o sea "0"
    - ORACIÓN esté vacía o sea "0"

    (Se aplica sobre el DataFrame de resultados detallados).
    """
    if df is None or df.empty:
        return df

    df = df.copy()
    mask_drop = pd.Series(False, index=df.index)

    if "Sugerencias" in df.columns:
        mask_drop |= df["Sugerencias"].apply(_is_empty_or_zero_cell)

    if "Oración" in df.columns:
        mask_drop |= df["Oración"].apply(_is_empty_or_zero_cell)

    if mask_drop.any():
        before = len(df)
        df = df.loc[~mask_drop].copy()
        logger.info(
            f"Filtro filas sin Sugerencias/Oración: "
            f"{before - len(df)} filas eliminadas; {len(df)} filas restantes."
        )

    return df


def to_excel_bytes(resultados_df: pd.DataFrame, resumen_completo_df: pd.DataFrame) -> bytes:
    # 🔹 OPCIONAL: asegurar limpieza antes de exportar
    resultados_df = _filter_resultados_empty_suggest_or_sentence(resultados_df)
    
    # Sanitizar dataframes de entrada para evitar caracteres ilegales en XML
    resultados_df = _sanitize_excel_df(resultados_df)
    resumen_completo_df = _sanitize_excel_df(resumen_completo_df)

    out = BytesIO()
    with pd.ExcelWriter(out, engine="xlsxwriter") as w:
        # -----------------------------
        # Hoja 1: Resultados
        # -----------------------------
        if resultados_df is None or resultados_df.empty:
            tmp = pd.DataFrame(columns=[
                "Archivo", "Página/Diapositiva", "BloqueTipo", "Mensaje",
                "Sugerencias", "Oración", "Contexto", "Regla", "Categoría"
            ])
            tmp = _sanitize_excel_df(tmp)
            tmp.to_excel(w, index=False, sheet_name="Resultados")
            ws = w.sheets["Resultados"]
            ws.set_column(0, 0, 40)
        else:
            resultados_df.to_excel(w, index=False, sheet_name="Resultados")
            ws = w.sheets["Resultados"]
            for i, col in enumerate(resultados_df.columns):
                try:
                    width = min(
                        60,
                        max(
                            12,
                            int(resultados_df[col].astype(str).str.len().quantile(0.9)) + 2
                        )
                    )
                except Exception:
                    width = 22
                ws.set_column(i, i, width)

        # -----------------------------
        # Hoja 2: ResumenIncidencias
        # -----------------------------
        if resultados_df is None or resultados_df.empty:
            resumen_inc = pd.DataFrame(columns=["Archivo", "TotalIncidencias"])
        else:
            resumen_inc = (
                resultados_df.groupby("Archivo")
                .size()
                .reset_index(name="TotalIncidencias")
                .sort_values("TotalIncidencias", ascending=False)
            )

        resumen_inc = _sanitize_excel_df(resumen_inc)
        resumen_inc.to_excel(w, index=False, sheet_name="ResumenIncidencias")
        ws2 = w.sheets["ResumenIncidencias"]
        ws2.set_column(0, 0, 40)
        ws2.set_column(1, 1, 22)

        # -----------------------------
        # Hoja 3: ResumenCompleto
        # -----------------------------
        if resumen_completo_df is None:
            resumen_completo_df = pd.DataFrame()

        resumen_completo_df = _sanitize_excel_df(resumen_completo_df)
        if resumen_completo_df is None:
            resumen_completo_df = pd.DataFrame()

        resumen_completo_df.to_excel(w, index=False, sheet_name="ResumenCompleto")
        ws3 = w.sheets["ResumenCompleto"]
        for i, col in enumerate(resumen_completo_df.columns):
            ws3.set_column(i, i, 28 if i > 0 else 50)

    # ⚠️ Importante: después del with YA NO se vuelve a escribir nada en `w`
    out.seek(0)
    return out.getvalue()

# ======================================================
# UTILIDADES DE RED (para descarga masiva)
# ======================================================
def _requests_available_or_warn() -> bool:
    """Verifica que la librería `requests` esté instalada."""
    if requests is None:
        st.error("Falta la librería `requests`. Instala con: `pip install requests`")
        return False
    return True

# ======================================================
# PDF to Word Transform
# ======================================================
BIB_HEAD_PATTERNS = [
    "fuentes bibliograficas",
    "referencias bibliograficas",
    "bibliografia",
    "referencias",
    "obras citadas",
]

class PDFBatchProcessor:
    def __init__(self, max_workers: int = None):
        self.max_workers = max_workers or min(4, os.cpu_count() or 1)
        self.cancel_event = threading.Event()
        self.progress_queue: "queue.Queue[Dict[str, Any]]" = queue.Queue()

    def process_single_pdf(self, pdf_path: str, output_dir: str, options: Dict) -> Dict:
        if self.cancel_event.is_set():
            return {"status": "cancelled", "file": pdf_path}
        if fitz is None or DocxDocument is None:
            return {
                "status": "error",
                "file": pdf_path,
                "error": "Faltan dependencias `pymupdf` o `python-docx`.",
                "elapsed_time": 0,
                "success": False,
            }
        try:
            start_time = datetime.now()
            pdf_name = Path(pdf_path).name
            self.progress_queue.put({"type": "file_start", "file": pdf_name, "total_pages": 0})
            success, output_path, stats = self._process_pdf_internal(pdf_path, output_dir, options)
            elapsed = (datetime.now() - start_time).total_seconds()
            result: Dict[str, Any] = {
                "status": "success" if success else "error",
                "file": pdf_path,
                "output": output_path,
                "stats": stats,
                "elapsed_time": elapsed,
                "success": success,
            }
            if not success:
                result["error"] = stats.get("error", "Error desconocido")
            return result
        except Exception as e:
            logger.error(f"Error procesando {pdf_path}: {e}")
            return {
                "status": "error",
                "file": pdf_path,
                "error": str(e),
                "elapsed_time": 0,
                "success": False,
            }

    def _process_pdf_internal(self, pdf_path: str, output_dir: str, options: Dict) -> Tuple[bool, str, Dict]:
        use_multithread = options.get("usar_multihilo", True)
        try:
            with fitz.open(pdf_path) as doc:
                num_pages = len(doc)
            self.progress_queue.put({"type": "file_pages", "file": Path(pdf_path).name, "total_pages": num_pages})
            output_path = Path(output_dir) / f"{Path(pdf_path).stem}.docx"
            doc_word = DocxDocument()
            if use_multithread and num_pages > 3:
                from concurrent.futures import ThreadPoolExecutor
                futures: List[Tuple[int, Any]] = []
                results: List[Tuple[int, str]] = []
                with ThreadPoolExecutor(max_workers=options.get("max_workers", 4)) as executor:
                    for page_num in range(num_pages):
                        if self.cancel_event.is_set():
                            break
                        future = executor.submit(self._extract_and_process_page, pdf_path, page_num, options)
                        futures.append((page_num, future))
                    for page_num, future in futures:
                        if self.cancel_event.is_set():
                            break
                        page_idx, text = future.result()
                        results.append((page_idx, text))
                        self.progress_queue.put({
                            "type": "page_progress",
                            "file": Path(pdf_path).name,
                            "page": page_num + 1,
                            "total": num_pages,
                        })
                results.sort(key=lambda x: x[0])
                page_texts = [text for _, text in results]
            else:
                page_texts: List[str] = []
                for page_num in range(num_pages):
                    if self.cancel_event.is_set():
                        break
                    _, text = self._extract_and_process_page(pdf_path, page_num, options)
                    page_texts.append(text)
                    self.progress_queue.put({
                        "type": "page_progress",
                        "file": Path(pdf_path).name,
                        "page": page_num + 1,
                        "total": num_pages,
                    })
            if self.cancel_event.is_set():
                return False, str(output_path), {"status": "cancelled"}
            for idx, text in enumerate(page_texts):
                if text.strip():
                    for line in text.split("\n"):
                        if line.strip():
                            doc_word.add_paragraph(line)
                if idx < len(page_texts) - 1:
                    doc_word.add_page_break()
            doc_word.save(str(output_path))
            stats = {
                "archivo": pdf_path,
                "nombre_archivo": Path(pdf_path).stem,
                "paginas_procesadas": num_pages,
                "archivo_salida": str(output_path),
                "tamano_salida": os.path.getsize(output_path) if os.path.exists(output_path) else 0,
                "errores": sum(1 for text in page_texts if "ERROR:" in text),
                "timestamp": datetime.now().isoformat(),
            }
            return True, str(output_path), stats
        except Exception as e:
            logger.error(f"Error interno procesando {pdf_path}: {e}")
            return False, str(Path(output_dir) / f"{Path(pdf_path).stem}.docx"), {"error": str(e)}

    def _extract_and_process_page(self, pdf_path: str, page_num: int, options: Dict) -> Tuple[int, str]:
        try:
            with fitz.open(pdf_path) as doc:
                if page_num >= len(doc):
                    return page_num, f"=== PÁGINA {page_num + 1} ===\nERROR: Página no existe"
                page = doc[page_num]
                raw_text = page.get_text()
                if len(raw_text.strip()) < 50:
                    raw_text = page.get_text("text")
                cleaned_text = self._clean_text(raw_text)
                if options.get("filtrar_bibliografia", False):
                    text_base = self._filter_references(cleaned_text)
                else:
                    text_base = cleaned_text
                text_no_formulas = self._filter_formulas(text_base)
                reformatted_text = self._reformat_sentences(text_no_formulas)
                result = f"=== PÁGINA {page_num + 1} ===\n{reformatted_text.strip()}"
                return page_num, result
        except Exception as e:
            error_msg = f"ERROR procesando página {page_num + 1}: {str(e)}"
            logger.error(f"{error_msg} en {pdf_path}")
            return page_num, f"=== PÁGINA {page_num + 1} ===\n{error_msg}"

    def _clean_text(self, text: str) -> str:
        replacements = {
            "\x00": "", "\x0c": "\n", "\uf0b7": "•", "\uf0a7": "§",
            "\uf0d8": "°", "\xad": "", "\t": "    ",
        }
        for old, new in replacements.items():
            text = text.replace(old, new)
        lines = [line.strip() for line in text.split("\n")]
        lines = [line for line in lines if line]
        final_lines: List[str] = []
        buffer: List[str] = []
        def flush_buffer():
            nonlocal buffer, final_lines
            if buffer:
                final_lines.append(" ".join(buffer))
                buffer = []
        for line in lines:
            if "=" in line:
                flush_buffer()
                final_lines.append(line)
                continue
            if len(line) < 80 and not line.endswith((".", "!", "?", ":", ";", ",", ")")):
                buffer.append(line)
            else:
                if buffer:
                    final_lines.append(" ".join(buffer + [line]))
                    buffer = []
                else:
                    final_lines.append(line)
        flush_buffer()
        return "\n".join(final_lines).strip()

    def _is_reference_line(self, line: str) -> bool:
        if not line or len(line.strip()) < 5:
            return False
        text = line.strip()
        if re.search(r"\(\d{4}[a-z]?\)", text) or re.search(r"\b\d{4}\.", text):
            return True
        if re.search(r"https?://\S+", text, re.IGNORECASE):
            return True
        numbers = sum(ch.isdigit() for ch in text)
        if numbers >= 4 and ("," in text or "." in text):
            if "pp." in text.lower() or "p." in text.lower():
                return True
        return False

    def _filter_references(self, text: str) -> str:
        if not text.strip():
            return text
        lines = text.split("\n")
        result: List[str] = []
        in_ref_block = False
        for i, line in enumerate(lines):
            norm = self._normalize_text(line)
            if in_ref_block:
                if not line.strip():
                    continue
                if self._is_reference_line(line):
                    continue
                in_ref_block = False
            is_bib_header = any(pattern in norm for pattern in BIB_HEAD_PATTERNS)
            if is_bib_header and self._is_reference_header(lines, i):
                in_ref_block = True
                continue
            result.append(line)
        return "\n".join(result).strip()

    def _filter_formulas(self, text: str) -> str:
        return text

    def _reformat_sentences(self, text: str) -> str:
        if not text.strip():
            return text
        text = re.sub(r"\s*\n\s*", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        text = re.sub(r"(\d+)\.\s+(\d{1,3})", r"\1.\2", text)
        text = re.sub(r"\b(\d+)\.\s+(?=[A-ZÁÉÍÓÚÑ])", r"\1§ ", text)
        text = re.sub(r"(?<!\d)\.\s+(?!\d)", ".\n", text)
        text = re.sub(r"\)\s+(?=[A-ZÁÉÍÓÚÑ¿])", ")\n", text)
        text = text.replace("§", ".")
        lines = [line.strip() for line in text.split("\n")]
        lines = [line for line in lines if line]
        return "\n".join(lines)

    def _normalize_text(self, text: str) -> str:
        if not text:
            return ""
        text = text.lower()
        text = "".join(c for c in unicodedata.normalize("NFD", text) if unicodedata.category(c) != "Mn")
        text = re.sub(r"\s+", " ", text).strip()
        return text

    def _is_reference_header(self, lines: List[str], idx: int) -> bool:
        if idx < 0 or idx >= len(lines):
            return False
        start = idx + 1
        end = min(len(lines), idx + 11)
        window = [line for line in lines[start:end] if line.strip()]
        if not window:
            return False
        total = len(window)
        ref_like = sum(1 for line in window if self._is_reference_line(line))
        return ref_like >= 2 or ref_like / total >= 0.4

# ======================================================
# DESCARGA MASIVA
# ======================================================
def _format_hms(seconds: float) -> str:
    seconds_int = int(max(0, seconds))
    m, s = divmod(seconds_int, 60)
    h, m = divmod(m, 60)
    if h > 0:
        return f"{h:02d}:{m:02d}:{s:02d}"
    return f"{m:02d}:{s:02d}"

def nombre_archivo_seguro(url: str, carpeta_destino: str, max_ruta: int = 240) -> str:
    nombre = url.split('/')[-1]
    nombre = nombre.split('?')[0]
    from urllib.parse import unquote as _unq
    nombre = _unq(nombre)
    caracteres_invalidos = '<>:"/\\|?*'
    for c in caracteres_invalidos:
        nombre = nombre.replace(c, '_')
    if not nombre.strip():
        nombre = "archivo_descargado"
    ruta_base = os.path.join(carpeta_destino, "")
    espacio_disponible = max_ruta - len(ruta_base)
    if espacio_disponible < 50:
        espacio_disponible = 50
    if len(nombre) > espacio_disponible:
        base, ext = os.path.splitext(nombre)
        if len(ext) > 10:
            ext = ext[:10]
        max_base = espacio_disponible - len(ext)
        if max_base < 1:
            max_base = 1
        base = base[:max_base]
        nombre = base + ext
    return nombre

def _ensure_unique_path(base_dir: Path, filename: str) -> Path:
    base_dir = Path(base_dir)
    base_dir.mkdir(parents=True, exist_ok=True)
    candidate = base_dir / filename
    if not candidate.exists():
        return candidate
    stem = Path(filename).stem
    ext = Path(filename).suffix
    counter = 1
    while True:
        alt = base_dir / f"{stem}_{counter}{ext}"
        if not alt.exists():
            return alt
        counter += 1

def _group_files_by_size(file_paths: List[str], max_mb: int = MAX_ZIP_BLOCK_MB) -> List[List[str]]:
    groups: List[List[str]] = []
    current_group: List[str] = []
    current_size = 0
    max_bytes = max_mb * 1024 * 1024
    for path in file_paths:
        try:
            size = os.path.getsize(path)
        except OSError:
            continue
        if size >= max_bytes:
            if current_group:
                groups.append(current_group)
                current_group = []
                current_size = 0
            groups.append([path])
            continue
        if current_group and current_size + size > max_bytes:
            groups.append(current_group)
            current_group = [path]
            current_size = size
        else:
            current_group.append(path)
            current_size += size
    if current_group:
        groups.append(current_group)
    return groups

def _zip_paths_to_bytes(file_paths: List[str]) -> bytes:
    buf = BytesIO()   # ✅ usar la clase ya importada
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in file_paths:
            if not os.path.exists(p):
                continue
            arcname = Path(p).name
            zf.write(p, arcname=arcname)
    buf.seek(0)
    return buf.getvalue()


def _run_descarga_masiva_streamlit(
    urls_archivos: List[str],
    progress_placeholder,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]], str, Optional[str]]:

    if not urls_archivos:
        tmp_dir = Path(tempfile.mkdtemp(prefix="utp_descarga_masiva_"))
        return [], [], str(tmp_dir), None
    tmp_dir = Path(tempfile.mkdtemp(prefix="utp_descarga_masiva_"))
    session = requests.Session()
    resultados: List[Dict[str, Any]] = []
    fallidos: List[Dict[str, Any]] = []
    total = len(urls_archivos)
    start_time = time.time()

    if progress_placeholder is not None and total > 0:
        render_task_progress(
            progress_placeholder,
            "Descargando archivos",
            0.0,
            f"0/{total} archivos",
        )

    for idx, url in enumerate(urls_archivos, start=1):
        url = str(url).strip()
        descargado_ok = False
        ultimo_error = ""
        nombre_archivo = ""
        ruta_archivo = None
        try:
            nombre_archivo = nombre_archivo_seguro(url, str(tmp_dir))
            ruta_archivo = tmp_dir / nombre_archivo
            if ruta_archivo.exists():
                base, ext = os.path.splitext(nombre_archivo)
                contador = 1
                while ruta_archivo.exists():
                    sufijo = f"_{contador}"
                    espacio_disponible = 240 - len(os.path.join(str(tmp_dir), ""))
                    max_base = espacio_disponible - len(ext) - len(sufijo)
                    if max_base < 1:
                        max_base = 1
                    base_trunc = base[:max_base]
                    nombre_archivo_alt = f"{base_trunc}{sufijo}{ext}"
                    ruta_archivo = tmp_dir / nombre_archivo_alt
                    contador += 1
            for intento in range(1, MAX_INTENTOS_DESCARGA + 1):
                try:
                    resp = session.get(url, stream=True, timeout=30, headers=REQUEST_HEADERS, allow_redirects=True)
                    if resp.status_code == 200:
                        tam_header = resp.headers.get("content-length")
                        if tam_header:
                            try:
                                tam_esperado = int(tam_header)
                            except ValueError:
                                tam_esperado = None
                        else:
                            tam_esperado = None
                        with open(ruta_archivo, "wb") as f:
                            for chunk in resp.iter_content(chunk_size=CHUNK_SIZE):
                                if chunk:
                                    f.write(chunk)
                        tam_real = os.path.getsize(ruta_archivo)
                        if (tam_esperado is not None and tam_esperado == 0) or tam_real == 0:
                            try:
                                os.remove(ruta_archivo)
                            except OSError:
                                pass
                            raise Exception(f"Archivo descargado con tamaño 0 (tam_esperado={tam_esperado}, tam_real={tam_real}).")
                        descargado_ok = True
                        resultados.append({
                            "url": url,
                            "nombre_archivo": ruta_archivo.name,
                            "ruta_archivo": str(ruta_archivo),
                            "status": "OK",
                        })
                        break
                    else:
                        raise Exception(f"Código HTTP: {resp.status_code}")
                except Exception as e:
                    ultimo_error = str(e)
                    logger.warning(f"Error en intento {intento} para {nombre_archivo}: {e}")
                    if ruta_archivo and ruta_archivo.exists():
                        try:
                            os.remove(ruta_archivo)
                        except OSError:
                            pass
                    if intento < MAX_INTENTOS_DESCARGA:
                        espera = min(60, 2 ** intento)
                        logger.info(f"Reintentando en {espera} segundos para {url}...")
                        time.sleep(espera)
            if not descargado_ok:
                fallidos.append({
                    "url": url,
                    "nombre_archivo": nombre_archivo,
                    "error": ultimo_error or "Error desconocido",
                })
        except Exception as e:
            logger.error(f"Error al procesar {url}: {e}")
            fallidos.append({
                "url": url,
                "nombre_archivo": nombre_archivo,
                "error": str(e),
            })
        elapsed = time.time() - start_time
        processed = idx
        pct = processed / total
        speed = processed / elapsed if elapsed > 0 else 0.0
        eta = (total - processed) / speed if speed > 0 else 0.0
        if progress_placeholder is not None and total > 0:
            detail = (
                f"{pct*100:.1f}% | {processed}/{total} archivos "
                f"[{_format_hms(elapsed)}<{_format_hms(eta)}, {speed:.2f} archivo/s]"
            )
            render_task_progress(
                progress_placeholder,
                "Descargando archivos",
                pct,
                detail,
            )

    csv_fallidos_path: Optional[Path] = None
    if fallidos:
        csv_fallidos_path = tmp_dir / "descargas_fallidas.csv"
        pd.DataFrame(fallidos).to_csv(csv_fallidos_path, index=False, encoding="utf-8-sig")
    return resultados, fallidos, str(tmp_dir), str(csv_fallidos_path) if csv_fallidos_path else None

# ======================================================
# UI HELPERS
# ======================================================
def apply_global_styles():
    st.markdown(
        """
        <style>
        /* Reducir tamaño base de fuente para que todo se vea más compacto */
        html {
            font-size: 11px;  /* en lugar de 16px */
        }

        [data-testid="stAppViewContainer"] { background: #f3f4f6; }
        [data-testid="stSidebar"] { background: #f9fafb; }

        .utp-hero {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 18px;
            padding: 1.8rem 2.4rem;
            color: #ffffff;
            margin-bottom: 1.8rem;
            box-shadow: 0 18px 40px rgba(76, 81, 191, 0.35);
            display: flex;
            align-items: center;
            gap: 1.0rem;
        }

        .utp-hero-icon {
            width: 3.1rem;
            height: 3.1rem;
            border-radius: 999px;
            background: rgba(255,255,255,0.18);
            display: flex;
            align-items: center;
            justify-content: center;
            font-size: 2.0rem;
        }
        .utp-hero-title { font-weight: 800; font-size: 1.8rem; margin-bottom: 0.15rem; }
        .utp-hero-sub { font-size: 0.92rem; opacity: 0.96; line-height: 1.4; }
        .utp-sidebar-brand {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            border-radius: 18px;
            padding: 1.0rem 1.1rem;
            color: #ffffff;
            box-shadow: 0 14px 32px rgba(76, 81, 191, 0.35);
            margin-bottom: 1.1rem;
        }
        .utp-sidebar-brand-title {
            font-weight: 800;
            font-size: 1.05rem;
            margin-bottom: 0.2rem;
            display: flex;
            align-items: center;
            gap: 0.4rem;
        }
        .utp-sidebar-brand-subtitle { font-size: 0.82rem; opacity: 0.92; }
        .utp-card {
            border-radius: 14px;
            border: 1px solid #e5e7eb;
            padding: 1.1rem 1.3rem 1.15rem 1.3rem;
            margin-bottom: 1.0rem;
            background: #ffffff;
            box-shadow: 0 10px 25px rgba(15,23,42,0.05);
        }
        .utp-step-row {
            display: flex;
            align-items: center;
            justify-content: space-between;
            margin-bottom: 0.7rem;
        }
        .utp-step-main {
            display: flex;
            align-items: center;
            gap: 0.55rem;
            font-size: 1.0rem;
            font-weight: 700;
            color: #111827;
        }
        .utp-step-number {
            display: inline-flex;
            align-items: center;
            justify-content: center;
            width: 26px;
            height: 26px;
            border-radius: 999px;
            background: #4f46e5;
            color: #ffffff;
            font-size: 0.9rem;
            font-weight: 700;
            box-shadow: 0 3px 8px rgba(79,70,229,0.45);
        }
        .utp-step-status {
            padding: 0.18rem 0.7rem;
            border-radius: 999px;
            font-size: 0.78rem;
            font-weight: 600;
            border: 1px solid transparent;
            white-space: nowrap;
        }
        .utp-step-status-ok { background-color: #dcfce7; color: #166534; border-color: #bbf7d0; }
        .utp-step-status-warn { background-color: #ffedd5; color: #9a3412; border-color: #fed7aa; }
        .utp-step-status-error { background-color: #fee2e2; color: #b91c1c; border-color: #fecaca; }
        .utp-step-header-simple {
            display: flex;
            align-items: center;
            gap: 0.55rem;
            font-size: 1.0rem;
            font-weight: 700;
            color: #111827;
            margin-bottom: 0.7rem;
        }
        .utp-step-header-simple .utp-step-number {
            width: 26px;
            height: 26px;
            border-radius: 999px;
            background: #4f46e5;
            color: #ffffff;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            font-size: 0.9rem;
            font-weight: 700;
        }
        .stButton>button {
            border-radius: 999px;
            font-weight: 700;
            padding: 0.6rem 1.3rem;
            border: none;
            transition: all 0.2s ease;
        }
        .stButton>button:hover {
            transform: translateY(-1px);
            box-shadow: 0 10px 25px rgba(79,70,229,0.45);
        }
        .stDataFrame { border-radius: 10px; border: 1px solid #e5e7eb; }
        .hero-reset-anchor + div[data-testid="stButton"] {
            margin-top: -3.0rem;
            margin-bottom: 0.6rem;
            display: flex;
            justify-content: flex-end;
            padding-right: 2.4rem;
        }
        .hero-reset-anchor + div[data-testid="stButton"] > button {
            position: relative;
            width: 44px;
            height: 44px;
            border-radius: 999px;
            padding: 0;
            background-color: #ff1654;
            color: transparent;
            box-shadow: 0 12px 28px rgba(15,23,42,0.45);
        }
        .hero-reset-anchor + div[data-testid="stButton"] > button::before {
            content: "↻";
            position: absolute;
            top: 50%;
            left: 50%;
            transform: translate(-50%, -50%);
            font-size: 22px;
            color: #ffffff;
        }
        .hero-reset-anchor + div[data-testid="stButton"] > button:hover {
            transform: translateY(-1px);
            box-shadow: 0 16px 32px rgba(15,23,42,0.55);
        }
        .utp-success-chip {
            display: inline-flex;
            align-items: center;
            justify-content: center;
            padding: 0.5rem 1.1rem;
            border-radius: 999px;
            background-color: #22c55e;
            color: #ffffff;
            font-weight: 700;
            font-size: 0.9rem;
            box-shadow: 0 10px 24px rgba(22,163,74,0.45);
            border: none;
            margin-top: 0.45rem;
        }

        .progress-bar-ui-task {
            margin-top: 0.35rem;
            margin-bottom: 0.35rem;
        }

        .progress-bar-ui-task-header {
            font-size: 0.85rem;
            font-weight: 600;
            color: #111827;
            margin-bottom: 0.25rem;
        }

        /* Contenedor de la barra: pill ancho tipo screenshot */
        .progress-bar-ui-task-track {
            position: relative;
            width: 100%;
            height: 20px;                      /* ⬅️ antes 18px: barra mucho más gruesa */
            border-radius: 999px;
            background: #e0f2ff;               /* azul muy claro como el fondo de la captura */
            overflow: hidden;
        }

        /* Barra de progreso azul */
        .progress-bar-ui-task-bar {
            position: absolute;
            top: 0;
            left: 0;
            height: 100%;                      /* ocupa todo el alto del track */
            display: flex;
            align-items: center;               /* centra el texto verticalmente */
            justify-content: flex-start;       /* texto pegado a la izquierda */
            padding: 0 12px;
            border-radius: 999px;
            background: linear-gradient(90deg, #1d8cf8, #1554d1);
            color: #ffffff;
            font-size: 0.80rem;
            font-weight: 700;
            text-shadow: 0 1px 2px rgba(15,23,42,0.4);
            min-width: 40px;                   /* asegura una “pastilla” azul visible */
            max-width: 100%;
            box-sizing: border-box;
        }

        .progress-bar-ui-task-subtext {
            margin-top: 0.25rem;
            font-size: 0.80rem;
            color: #4b5563;
        }


        </style>
    """, unsafe_allow_html=True)

def render_sidebar_header():
    st.markdown("""
        <div class="utp-sidebar-brand">
            <div class="utp-sidebar-brand-title">
                <span>📚</span><span>Plataforma - GrammarScan</span>
            </div>
            <div class="utp-sidebar-brand-subtitle">
                Revisión automatizada de ortografía y gramática de diversos documentos académicos.
            </div>
        </div>
    """, unsafe_allow_html=True)


def render_hero(title: str, subtitle: str, icon: str = "📚"):
    st.markdown(f"""
        <div class="utp-hero">
            <div class="utp-hero-icon">{icon}</div>
            <div>
                <div class="utp-hero-title">{title}</div>
                <div class="utp-hero-sub">{subtitle}</div>
            </div>
        </div>
    """, unsafe_allow_html=True)

def render_step_header_html(step_label: str, title: str, status: str) -> str:
    map_text = {"ok": "Listo", "warn": "Pendiente", "error": "Falta"}
    map_class = {"ok": "utp-step-status-ok", "warn": "utp-step-status-warn", "error": "utp-step-status-error"}
    status_text = map_text.get(status, "Pendiente")
    status_class = map_class.get(status, "utp-step-status-warn")
    return f"""
    <div class="utp-step-row">
        <div class="utp-step-main">
            <span class="utp-step-number">{step_label}</span>
            <span>{title}</span>
        </div>
        <div class="utp-step-status {status_class}">{status_text}</div>
    </div>
    """

def render_simple_step_header(step_label: str, title: str):
    st.markdown(f"""
        <div class="utp-step-header-simple">
            <span class="utp-step-number">{step_label}</span>
            <span>{title}</span>
        </div>
    """, unsafe_allow_html=True)

def ui_card_open():
    st.markdown('<div class="utp-card">', unsafe_allow_html=True)

def ui_card_close():
    st.markdown("</div>", unsafe_allow_html=True)

def render_success_chip(text: str):
    st.markdown(f'<div class="utp-success-chip">{text}</div>', unsafe_allow_html=True)

def render_task_progress(
    placeholder,
    title: str,
    pct: float,
    detail: str = "",
) -> None:
    """Barra de progreso unificada tipo progress-bar-ui-task."""
    pct = max(0.0, min(1.0, float(pct or 0.0)))
    pct_label = f"{pct * 100:.1f}%"
    detail_html = f'<div class="progress-bar-ui-task-subtext">{detail}</div>' if detail else ""
    html = f"""
    <div class="progress-bar-ui-task">
        <div class="progress-bar-ui-task-header">{title}</div>
        <div class="progress-bar-ui-task-track">
            <div class="progress-bar-ui-task-bar" style="width: {pct * 100:.1f}%;">
                {pct_label}
            </div>
        </div>
        {detail_html}
    </div>
    """
    placeholder.markdown(html, unsafe_allow_html=True)


# ======================================================
# ESTADO DE SESIÓN
# ======================================================
def init_session_state():
    if "module" not in st.session_state:
        st.session_state["module"] = "Home"
    if "module_radio" not in st.session_state:
        st.session_state["module_radio"] = st.session_state["module"]

    # ======================
    # Descarga masiva desde Excel
    # ======================
    st.session_state.setdefault("descarga_zip_bytes", None)
    st.session_state.setdefault("descarga_resultados", None)
    st.session_state.setdefault("descarga_fallidos", None)
    st.session_state.setdefault("descarga_download_dir", None)
    st.session_state.setdefault("descarga_fallidos_csv", None)

    # ======================
    # PDF → Word / extracción de texto
    # ======================
    st.session_state.setdefault("extraccion_zip_bytes", None)
    st.session_state.setdefault("extraccion_resultados", None)
    st.session_state.setdefault("extraccion_errores", None)

    # ======================
    # Pipeline global (descarga → PDF → Word → GrammarScan)
    # ======================
    st.session_state.setdefault("pipeline_pdf_signature", None)
    st.session_state.setdefault("pipeline_pdf_done", False)
    st.session_state.setdefault("pipeline_pdf_results", None)
    st.session_state.setdefault("pipeline_pdf_errors", None)

    st.session_state.setdefault("pipeline_docx_paths", [])
    st.session_state.setdefault("pipeline_docx_meta", {})
    st.session_state.setdefault("pipeline_pptx_paths", [])
    st.session_state.setdefault("pipeline_pptx_meta", {})
    st.session_state.setdefault("pipeline_word_inputs_count", 0)
    st.session_state.setdefault("pipeline_ppt_inputs_count", 0)

    st.session_state.setdefault("pipeline_reset_token", 0)

    # ======================
    # Bulk download desde Excel
    # ======================
    st.session_state.setdefault("bulk_has_valid_urls", False)
    st.session_state.setdefault("bulk_urls_archivos", None)
    st.session_state.setdefault("bulk_excel_df", None)
    st.session_state.setdefault("bulk_url_mapping", None)
    st.session_state.setdefault("pipeline_bulk_signature", None)
    st.session_state.setdefault("pipeline_bulk_done", False)

    # Directorio temporal para archivos subidos manualmente
    st.session_state.setdefault("pipeline_manual_dir", None)

    # ======================
    # GrammarScan
    # ======================
    st.session_state.setdefault("gs_uploader_key", 0)
    st.session_state.setdefault("gs_lang", "es")
    st.session_state.setdefault("gs_max_chars", 30000)
    st.session_state.setdefault("gs_workers", 4)
    st.session_state.setdefault("gs_excluir_biblio", True)
    st.session_state.setdefault("gs_modismos", False)
    st.session_state.setdefault("gs_final_df", None)
    st.session_state.setdefault("gs_resumen_completo_df", None)
    st.session_state.setdefault("gs_metrics", None)
    st.session_state.setdefault("gs_elapsed", 0.0)
    st.session_state.setdefault("gs_last_files_signature", None)

    st.session_state.setdefault("gs_excel_bytes", None)
    st.session_state.setdefault("gs_excel_autotrigger_done", False)

def reset_report_broken_pipeline():
    """
    Resetea el pipeline de:
    - Descarga masiva desde Excel
    - Procesamiento PDF → Word
    - Rutas de documentos para GrammarScan

    El estado propio de GrammarScan se limpia aparte en `reset_grammarscan_state()`.
    """
    keys_to_clear = [
        # Descarga masiva desde Excel
        "pipeline_bulk_signature",
        "pipeline_bulk_done",
        "bulk_has_valid_urls",
        "bulk_urls_archivos",
        "bulk_excel_df",
        "bulk_url_mapping",

        # Resultados de descarga
        "descarga_zip_bytes",
        "descarga_resultados",
        "descarga_fallidos",
        "descarga_download_dir",
        "descarga_fallidos_csv",

        # Procesamiento PDF → Word
        "pipeline_pdf_signature",
        "pipeline_pdf_done",
        "pipeline_pdf_results",
        "pipeline_pdf_errors",
        "extraccion_resultados",
        "extraccion_errores",
        "extraccion_zip_bytes",
        "extr_usar_multihilo",
        "extr_max_workers",

        # Rutas de documentos ya normalizados (para GrammarScan)
        "pipeline_docx_paths",
        "pipeline_docx_meta",
        "pipeline_pptx_paths",
        "pipeline_pptx_meta",
        "pipeline_word_inputs_count",
        "pipeline_ppt_inputs_count",

        # Directorio temporal de subidas manuales
        "pipeline_manual_dir",
    ]

    for k in keys_to_clear:
        st.session_state.pop(k, None)

    st.session_state["pipeline_reset_token"] = st.session_state.get("pipeline_reset_token", 0) + 1

def reset_grammarscan_state():
    keys_to_clear = [
        "gs_lang", "gs_max_chars", "gs_workers", "gs_excluir_biblio",
        "gs_modismos", "gs_final_df", "gs_resumen_completo_df", "gs_metrics",
        "gs_elapsed", "gs_last_files_signature",
        "gs_excel_bytes", "gs_excel_autotrigger_done",   # 👈 IMPORTANTE
    ]
    for k in keys_to_clear:
        if k in st.session_state:
            del st.session_state[k]

    st.session_state["gs_uploader_key"] += 1
    if "gs_uploader" in st.session_state:
        del st.session_state["gs_uploader"]

def reset_full_pipeline():
    reset_report_broken_pipeline()
    reset_grammarscan_state()
    st.rerun()

def on_change_module():
    st.session_state["module"] = st.session_state["module_radio"]

# ======================================================
# FUNCIONES AUXILIARES
# ======================================================
def _read_excel_safe(uploaded_file) -> pd.DataFrame:
    try:
        return pd.read_excel(uploaded_file)
    except Exception as e:
        raise RuntimeError(f"No se pudo leer el Excel: {e}") from e

def build_files_signature(uploaded_files) -> str:
    parts = []
    for uf in uploaded_files:
        size = getattr(uf, "size", None)
        parts.append(f"{uf.name}:{size}")
    parts.sort()
    return "|".join(parts)

def expand_uploaded_files(ups) -> List[LogicalFileSource]:
    logical_files: List[LogicalFileSource] = []

    for up in ups:
        name = up.name
        ext = os.path.splitext(name)[1].lower()

        if ext == ".zip":
            try:
                zip_bytes = up.getvalue()
                with zipfile.ZipFile(BytesIO(zip_bytes)) as zf:
                    for info in zf.infolist():
                        if info.is_dir():
                            continue
                        inner_ext = os.path.splitext(info.filename)[1].lower()
                        if inner_ext not in ALLOWED_DOC_EXTS:
                            continue

                        display_name = f"{Path(name).stem}/{Path(info.filename).name}"

                        def _make_reader(zbytes: bytes, inner_name: str) -> Callable[[], bytes]:
                            def _reader() -> bytes:
                                with zipfile.ZipFile(BytesIO(zbytes)) as _zf:
                                    return _zf.read(inner_name)
                            return _reader

                        logical_files.append(
                            LogicalFileSource(
                                display_name=display_name,
                                ext=inner_ext,
                                read_bytes=_make_reader(zip_bytes, info.filename),
                            )
                        )
            except Exception as e:
                st.error(f"Error leyendo ZIP '{name}': {e}")
                continue

        elif ext in ALLOWED_DOC_EXTS:
            logical_files.append(
                LogicalFileSource(
                    display_name=name,
                    ext=ext,
                    read_bytes=up.getvalue,
                )
            )
        else:
            st.warning(f"Archivo omitido por extensión no soportada: {name}")

    return logical_files


def process_grammarscan_files(
    ups,
    lang_code: str,
    max_chars_call: int,
    workers: int,
    excluir_biblio: bool,
    analizar_modismos: bool,
):
    try:
        _ = get_language_tool(lang_code)
    except Exception as e:
        st.error(f"No se pudo iniciar LanguageTool local: {e}")
        return pd.DataFrame([]), pd.DataFrame([]), {"total": 0, "n_inc": 0, "n_zero": 0, "n_err": 0}, 0.0

    modismos_patterns: List[ModismoPattern] = []
    if analizar_modismos and lang_code.startswith("es"):
        script_dir = os.path.dirname(os.path.abspath(__file__)) if "__file__" in globals() else os.getcwd()
        modismos_path = os.path.join(script_dir, "modismos_ar.xlsx")
        try:
            modismos_patterns = get_modismos_patterns(modismos_path)
            st.success(f"Diccionario de modismos cargado: {len(modismos_patterns)} entradas.")
        except Exception as e:
            st.error(f"No se pudieron cargar los modismos desde '{modismos_path}': {e}")
            modismos_patterns = []

    logical_files = expand_uploaded_files(ups)
    total_seleccionados = len(logical_files)

    if total_seleccionados == 0:
        return pd.DataFrame([]), pd.DataFrame([]), {"total": 0, "n_inc": 0, "n_zero": 0, "n_err": 0}, 0.0

    all_dfs: List[pd.DataFrame] = []
    resumen_rows: List[Dict[str, Any]] = []

    progress_task = st.empty()
    if total_seleccionados > 0:
        render_task_progress(
            progress_task,
            "Analizando documentos",
            0.0,
            f"0/{total_seleccionados} archivos",
        )

    t0 = time.time()

    for i, lf in enumerate(logical_files, start=1):
        pct = i / total_seleccionados if total_seleccionados > 0 else 1.0
        detail = f"{i}/{total_seleccionados} archivos procesados" if total_seleccionados > 0 else ""
        render_task_progress(
            progress_task,
            "Analizando documentos",
            pct,
            detail,
        )
        try:
            data = lf.read_bytes()
            ext = lf.ext
            df = analyze_file(
                lf.display_name,
                data,
                lang_code,
                max_chars_call,
                workers,
                excluir_bibliografia=excluir_biblio,
                modismos_patterns=modismos_patterns,
                analizar_modismos=analizar_modismos,
            )
            if not df.empty:
                all_dfs.append(df)
                resumen_rows.append({
                    "Archivo": lf.display_name,
                    "Extension": ext,
                    "Estado": "Con incidencias",
                    "TotalIncidencias": int(df.shape[0]),
                    "Detalle": ""
                })
            else:
                resumen_rows.append({
                    "Archivo": lf.display_name,
                    "Extension": ext,
                    "Estado": "Sin incidencias o sin texto",
                    "TotalIncidencias": 0,
                    "Detalle": ""
                })
        except Exception as e:
            resumen_rows.append({
                "Archivo": lf.display_name,
                "Extension": lf.ext,
                "Estado": "Error",
                "TotalIncidencias": None,
                "Detalle": safe_str(e)
            })
            st.error(f"Error procesando {lf.display_name}: {e}")

    resumen_completo_df = pd.DataFrame(resumen_rows)
    n_inc = int(resumen_completo_df.query("Estado == 'Con incidencias'")["Archivo"].nunique()) if not resumen_completo_df.empty else 0
    n_zero = int(resumen_completo_df.query("Estado == 'Sin incidencias o sin texto'")["Archivo"].nunique()) if not resumen_completo_df.empty else 0
    n_err = int(resumen_completo_df.query("Estado == 'Error'")["Archivo"].nunique()) if not resumen_completo_df.empty else 0

    if any(len(df) for df in all_dfs):
        final_df = pd.concat(all_dfs, ignore_index=True)
    else:
        final_df = pd.DataFrame([])

    # 🔹 NUEVO: limpiar filas sin Sugerencias u Oración (vacías o "0")
    final_df = _filter_resultados_empty_suggest_or_sentence(final_df)

    elapsed = time.time() - t0
    metrics = {
        "total": total_seleccionados,
        "n_inc": n_inc,
        "n_zero": n_zero,
        "n_err": n_err,
    }
    return final_df, resumen_completo_df, metrics, elapsed


# ======================================================
# PÁGINAS / MÓDULOS
# ======================================================

def page_home():
    # Hero principal
    render_hero(
        title=APP_TITLE,
        subtitle=(
            "Revisión automatizada y validación inteligente de ortografía y gramática "
            "en documentos académicos y administrativos."
        ),
        icon="📚",
    )


    # Card principal de contenido
    ui_card_open()

    # Sección: Home + propósito general
    st.markdown(
        """
        ### 🏠 Home

        UTP GrammarScan es una herramienta inteligente desarrollada para el análisis automatizado de información
        contenida en diversos tipos de archivos académicos y administrativos.
        Su objetivo es optimizar los procesos de revisión documental y reducir el tiempo de revisión manual, 
        asegurando la calidad lingüística y la consistencia de los textos académicos producidos dentro de la institución.

        ### 🎯 Propósito de la Plataforma

        Automatizar el proceso completo de revisión documental desde la recolección de materiales hasta el análisis 
        lingüístico avanzado, ofreciendo una solución integral que:

        - Unifica múltiples fuentes de documentos (descarga masiva, carga manual, procesamiento automático).
        - Estandariza formatos (PDF, Word, PPT → Texto estructurado).
        - Analiza contenido con motores de gramática y ortografía profesional.
        - Genera reportes detallados y exportables listos para revisión.
        """,
        unsafe_allow_html=False,
    )

    # Sección: Funcionalidades principales
    st.markdown(
        """
        ---
        ### ✨ Funcionalidades Principales

        - **Descarga Masiva Inteligente (Extracción Automática)**. Identifica y descarga automáticamente documentos **PDF, Word y PPT** desde listados de URLs en Excel.
        - **Transformación Avanzada de Documentos PDF a Word**. Convierte documentos PDF a formato Word manteniendo la estructura y contenido textual.  
        - **Análisis Lingüístico Avanzado**. Revisión ortográfica y gramatical con LanguageTool. Detección de modismos argentinos con diccionario personalizado. 
        - **Reporte Status Excel**. Genera reportes en Excel detallados (Archivo, página, error, sugerencia, contexto).  
        """,
        unsafe_allow_html=False,
    )

    # Sección: Flujo de trabajo (NUEVA VERSIÓN)
    st.markdown("---")
    st.markdown("### ✨ Flujo de Trabajo")

    # 1) Resumen en tarjetas horizontales CON FLECHAS
    pasos_resumen = [
        ("1", "Carga de URLs", "Excel con enlaces"),
        ("2", "Descarga masiva", "Documentos origen"),
        ("3", "Transformación", "PDF → Word"),
        ("4", "Extracción", "Análisis de Documentos"),
        ("5", "Validación", "Ortográfica y Gramatical"),
        ("6", "Reporte final", "Excel Status"),
    ]

    # [card, arrow, card, arrow, ..., card]
    pesos = []
    for i in range(len(pasos_resumen)):
        pesos.append(4)  # columna de tarjeta
        if i < len(pasos_resumen) - 1:
            pesos.append(1)  # columna de flecha

    cols = st.columns(pesos)

    idx = 0
    for i, (numero, titulo, subtitulo) in enumerate(pasos_resumen):
        # tarjeta
        with cols[idx]:
            st.markdown(
                f"""
                <div style="
                    background-color: #ffffff;
                    border-radius: 14px;
                    padding: 16px 18px;
                    box-shadow: 0 1px 3px rgba(15, 23, 42, 0.12);
                    text-align: center;
                    font-size: 0.85rem;
                ">
                    <div style="
                        display:inline-flex;
                        align-items:center;
                        justify-content:center;
                        width:32px;
                        height:32px;
                        border-radius:8px;
                        background:#3b82f6;
                        color:#ffffff;
                        font-weight:700;
                        margin-bottom:6px;
                    ">
                        {numero}
                    </div>
                    <div style="font-weight: 600;">{titulo}</div>
                    <div style="color: #6b7280; font-size: 0.75rem;">{subtitulo}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )

        idx += 1

        # flecha (entre tarjetas, excepto después de la última)
        if i < len(pasos_resumen) - 1:
            with cols[idx]:
                st.markdown(
                    """
                    <div style="
                        display:flex;
                        align-items:center;
                        justify-content:center;
                        min-height:120px;
                    ">
                        <span style="font-size:1.8rem; color:#9ca3af;">➜</span>
                    </div>
                    """,
                    unsafe_allow_html=True,
                )

            idx += 1

    st.markdown("")  # pequeño espacio debajo del timeline

 # 2) Detalle por paso en expanders (6 pasos completos)
    with st.expander("🔹 Paso 1: Carga de Excel de URLs", expanded=False):
        st.markdown(
            """
            - **Formato flexible**: soporta múltiples estructuras de columnas en Excel.  
            - **Normalización automática**: corrige y estandariza formatos de URLs.  
            - **Validación preliminar**: detecta problemas estructurales antes del procesamiento.
            """
        )

    with st.expander("🔹 Paso 2: Descarga Masiva", expanded=False):
        st.markdown(
            """
            - **Procesamiento automático**: descarga simultánea de múltiples documentos.  
            - **Gestión de errores**: registro detallado de fallos con motivos específicos.  
            - **Organización automática**: archivos descargados listos para procesamiento posterior.
            """
        )

    with st.expander("🔹 Paso 3: Transformación de Documentos", expanded=False):
        st.markdown(
            """
            - **Conversión** **PDF → Word**: extracción textual manteniendo referencias.  
            - **Procesamiento paralelo**: uso eficiente de recursos para documentos grandes.  
            - **Preservación de metadatos**: mantenimiento de información de origen.
            """
        )

    with st.expander("🔹 Paso 4: Análisis de Documentos", expanded=False):
        st.markdown(
            """
            - **Análisis exhaustivo**: procesamiento completo de documentos convertidos (PDF, Word, PPT) para extracción de contenido.  
            - **Contexto completo**: preparación optimizada del contenido para la validación lingüística.
            """
        )

    with st.expander("🔹 Paso 5: Ortografía y Gramática", expanded=False):
        st.markdown(
            """
            - **Verificación en tiempo real**: Análisis con LanguageTool.  
            - **Clasificación modismos**: detección de modismos argentinos específicos.  
            """
        )

    with st.expander("🔹 Paso 6: Reporte Final", expanded=False):
        st.markdown(
            """
            - **Excel estructurado**: exportación de reporte Excel.  
            - **Descarga Automatizada**: Descarga automática al completar el análisis.
            """
        )
    # Sección: Seguridad
    st.markdown(
        """
        ---
        ### 🔒 Seguridad y Privacidad

        - **Procesamiento local**: no se almacenan documentos en servidores externos.  
        - **Metadatos anónimos**: solo se registra información necesaria para el análisis.  
        - **Sin persistencia**: los archivos temporales se eliminan después del procesamiento.  
        - **Control total**: el usuario mantiene control completo sobre sus documentos.
        """,
        unsafe_allow_html=False,
    )

    ui_card_close()

def render_report_grammarscan():
    # ======================================================
    # 1. Bulk Document (PDF, WORD and PPT) Download
    # ======================================================
    render_hero(
        "Descarga Masiva Automática de Documentos (PDF's - Word y PPT) desde un Reporte Excel",
        "Sube tu Excel con URLs y descarga automáticamente todos los documentos y el reporte en Excel.",
        "⬇️",
    )
    
    # Botón Reiniciar
    st.markdown('<div class="hero-reset-anchor"></div>', unsafe_allow_html=True)
    if st.button("Reiniciar", key="btn_reset_report_broken"):
        reset_full_pipeline()

    # ------------------------------------------------------------------
    # VENTANA COLAPSABLE: Descarga masiva de documentos (PDF, Word, PPT)
    #   - Colapsada por defecto.
    #   - Al pulsar Reiniciar se incrementa pipeline_reset_token y se
    #     genera un label interno distinto añadiendo caracteres invisibles
    #     (\u200B), lo que fuerza a Streamlit a recrear el expander
    #     completamente colapsado, sin usar key.
    # ------------------------------------------------------------------

        # ------------------------------------------------------------------
    # NUEVA VENTANA: Filtrar bibliografía y modismos Arg.
    #   - Colapsada por defecto
    #   - Solo "Excluir bibliografía" viene activado por defecto
    # ------------------------------------------------------------------
    with st.expander("Filtrar bibliografia y modismos Arg.", expanded=False):
        st.markdown(
            "Configura qué contenido quieres excluir o analizar antes de iniciar el flujo.",
            unsafe_allow_html=False,
        )

        col_bib, col_mod = st.columns(2)

        with col_bib:
            st.checkbox(
                "Excluir secciones/entradas de bibliografía (APA, MLA, IEEE, Vancouver)",
                key="gs_excluir_biblio",
                value=st.session_state.get("gs_excluir_biblio", True),
            )

        with col_mod:
            st.checkbox(
                "Detectar modismos argentinos (modismos_ar.xlsx)",
                key="gs_modismos",
                value=st.session_state.get("gs_modismos", False),
            )

    exp_base_label = "Descarga masiva de documentos (PDF, Word, PPT)"
    reset_counter = st.session_state.get("pipeline_reset_token", 0)
    exp_label_internal = exp_base_label + ("\u200B" * reset_counter)

    with st.expander(exp_label_internal, expanded=False):

        # ---------- TARJETA: Selección de Excel (Paso 1) ----------
        ui_card_open()
        step1_ph = st.empty()
        
        bulk_uploader_key = f"pipeline_bulk_excel_uploader_{st.session_state.get('pipeline_reset_token', 0)}"
        uploaded_excel = st.file_uploader(
            "Seleccione el archivo Excel que contiene las URLs de los documentos a descargar",
            type=["xlsx", "xls"],
            key=bulk_uploader_key,
        )
        
        file_ok = uploaded_excel is not None
        step1_ph.markdown(
            render_step_header_html(
                "1",
                "Seleccione el archivo Excel que contiene las URLs de los documentos a descargar",
                "ok" if file_ok else "warn",
            ),
            unsafe_allow_html=True,
        )
        
        bulk_urls_archivos: List[str] = []
        df_in_bulk: Optional[pd.DataFrame] = None
        
        if file_ok:
            try:
                df_in_bulk = _read_excel_safe(uploaded_excel)
            except Exception as e:
                st.error(str(e))
            else:
                st.session_state["bulk_excel_df"] = df_in_bulk
                if "url" not in df_in_bulk.columns:
                    st.error("El Excel no contiene la columna requerida: **url**.")
                    st.caption(f"Columnas detectadas: {', '.join(map(str, df_in_bulk.columns.tolist()))}")
                else:
                    # Ya está normalizada, solo descartamos NaN
                    df_urls = df_in_bulk["url"].dropna()

                    # Aseguramos que lo que se manda a descarga también vaya limpio
                    bulk_urls_archivos = [
                        str(u).strip()
                        for u in df_urls
                        if str(u).strip().lower().endswith(DESC_EXT_PERMITIDAS)
                    ]
                    total_urls = len(df_urls)
                    total_permitidas = len(bulk_urls_archivos)

                    # 🔹 Límite duro de seguridad cuando se ejecuta en Streamlit Cloud
                    if IS_STREAMLIT_CLOUD and total_permitidas > MAX_BULK_URLS_CLOUD:
                        st.error(
                            f"El Excel contiene {total_permitidas} URLs descargables, "
                            f"lo cual supera el límite de {MAX_BULK_URLS_CLOUD} URLs "
                            "por ejecución en Streamlit Cloud. "
                            "Para más de 700 URLs, divide el Excel o ejecuta la herramienta en local."
                        )
                        st.info(
                            "📌 Documento demasiado grande para **Streamlit Cloud**; "
                            "divide el Excel en varios archivos más pequeños "
                            "o ejecuta la herramienta en tu equipo local."
                        )
                        # No habilitar la descarga masiva automática
                        st.session_state["bulk_has_valid_urls"] = False
                        st.session_state["bulk_urls_archivos"] = None
                    else:
                        if total_permitidas == 0:
                            st.warning(
                                "No se encontraron URLs que terminen en .ppt, .pptx, .pdf, .doc o .docx."
                            )
                        else:
                            st.session_state["bulk_has_valid_urls"] = True
                            st.session_state["bulk_urls_archivos"] = bulk_urls_archivos
                        
                        try:
                            excel_bytes = uploaded_excel.getbuffer()
                            bulk_signature = (uploaded_excel.name, len(excel_bytes))
                        except Exception:
                            bulk_signature = (uploaded_excel.name, 0)
                        
                        prev_bulk_sig = st.session_state.get("pipeline_bulk_signature")
                        if prev_bulk_sig != bulk_signature:
                            st.session_state["pipeline_bulk_signature"] = bulk_signature
                            st.session_state["pipeline_bulk_done"] = False
        else:
            st.caption("Carga un archivo Excel para continuar con la descarga masiva.")
        
        ui_card_close()
        
        # ---------- TARJETA: Procesar Descarga Masiva (Paso 2) ----------
        ui_card_open()
        render_simple_step_header("2", "Procesar Descarga Masiva")
        
        if not _requests_available_or_warn():
            ui_card_close()
            # salimos sólo del bloque del expander
        else:
            progress_task_bulk = st.empty()
            
            urls_archivos_state = st.session_state.get("bulk_urls_archivos") or []
            auto_trigger_bulk = bool(urls_archivos_state) and not st.session_state.get("pipeline_bulk_done", False)
            
            if urls_archivos_state and auto_trigger_bulk:
                try:
                    render_task_progress(
                        progress_task_bulk,
                        "Descargando archivos",
                        0.0,
                        "Preparando descarga masiva...",
                    )
                    with st.spinner("Descargando archivos..."):
                        resultados, fallidos, download_dir, csv_fallidos_path = _run_descarga_masiva_streamlit(
                            urls_archivos_state,
                            progress_placeholder=progress_task_bulk,
                        )
    
                    st.session_state["descarga_resultados"] = resultados
                    st.session_state["descarga_fallidos"] = fallidos
                    st.session_state["descarga_download_dir"] = download_dir
                    st.session_state["descarga_fallidos_csv"] = csv_fallidos_path
                    st.session_state["pipeline_bulk_done"] = True
    
                except Exception as e:
                    progress_task_bulk.empty()
                    st.error(f"Ocurrió un error durante la descarga masiva: {e}")
            
            resultados_ready = st.session_state.get("descarga_resultados") or []
            fallidos_ready = st.session_state.get("descarga_fallidos") or []
            
            if not (resultados_ready or fallidos_ready):
                if not urls_archivos_state:
                    st.caption("Primero carga un Excel válido con URLs para poder procesar la descarga.")
            
            if st.session_state.get("pipeline_bulk_done", False):
                render_success_chip("Descarga masiva completada")
            
            ui_card_close()
        
        # ---------- TARJETA: Descargar ZIP (Paso 3) ----------
        ui_card_open()
        render_simple_step_header("3", "Descargar todos los archivos (PDF, Word, PPT) (ZIP)")
        
        resultados_ready = st.session_state.get("descarga_resultados") or []
        download_dir = st.session_state.get("descarga_download_dir")
        csv_fallidos_path = st.session_state.get("descarga_fallidos_csv")
        
        file_paths: List[str] = []
        for r in resultados_ready:
            p = r.get("ruta_archivo")
            if p and os.path.exists(p):
                file_paths.append(str(p))
        
        if not resultados_ready:
            st.warning("Primero ejecuta el paso 2 para generar las descargas.")
        elif not file_paths:
            st.info(
                "No hay archivos descargados correctamente para comprimir en ZIP "
                "(es posible que todas las descargas hayan fallado)."
            )
        else:
            groups = _group_files_by_size(file_paths, max_mb=MAX_ZIP_BLOCK_MB)
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            if len(groups) == 1:
                zip_bytes = _zip_paths_to_bytes(groups[0])
                zip_name = f"Descarga_Masiva_Documentos_{ts}.zip"
                st.download_button(
                    "⬇️ Descargar todos los archivos (ZIP)",
                    data=zip_bytes,
                    file_name=zip_name,
                    mime="application/zip",
                )
            else:
                st.info(
                    f"Los archivos descargados se han dividido en **{len(groups)} ZIPs** "
                    f"para que cada uno tenga como máximo ~{MAX_ZIP_BLOCK_MB} MB."
                )
                for idx, group in enumerate(groups, start=1):
                    zip_bytes = _zip_paths_to_bytes(group)
                    zip_name = f"Descarga_Masiva_Documentos_{ts}_parte{idx}.zip"
                    st.download_button(
                        f"⬇️ Descargar ZIP parte {idx}",
                        data=zip_bytes,
                        file_name=zip_name,
                        mime="application/zip",
                        key=f"btn_zip_part_{idx}",
                    )
        
        if csv_fallidos_path and os.path.exists(csv_fallidos_path):
            try:
                with open(csv_fallidos_path, "rb") as fh:
                    csv_bytes = fh.read()
                st.download_button(
                    "⬇️ Descargar CSV de descargas fallidas",
                    data=csv_bytes,
                    file_name="descargas_fallidas.csv",
                    mime="text/csv",
                )
            except Exception as e:
                st.warning(f"No se pudo leer el CSV de descargas fallidas: {e}")
        
        ui_card_close()
    
    # ======================================================
    # 2. PDF, WORD and PPT to Word Transformation (ZIP)
    # ======================================================
    render_hero(
        "Carga Directa de Documentos (PDF's - Word y PPT) y ZIP",
        "Arrastra tus PDFs, Word, PPT o ZIP y el sistema los procesa automáticamente.",
        "🧲",
    )
    # ... (resto de la función sigue igual)

    
    if fitz is None or DocxDocument is None:
        ui_card_open()
        st.error(
            "Faltan dependencias para este módulo.\n\n"
            "- Instala `pymupdf` (fitz)\n"
            "- Instala `python-docx`\n\n"
            "Luego vuelve a desplegar la aplicación."
        )
        ui_card_close()
        return
    
    ui_card_open()
    step_pdf1 = st.empty()
    
    pdf_uploader_key = f"pipeline_pdf_uploader_ultra_{st.session_state.get('pipeline_reset_token', 0)}"
    uploaded_files = st.file_uploader(
        "Selecciona uno o más archivos PDF, Word, PPT o ZIP (ZIP con PDFs/Word/PPT en su interior)",
        type=["pdf", "docx", "pptx", "doc", "ppt", "zip"],
        accept_multiple_files=True,
        key=pdf_uploader_key,
        help="Si ya ejecutaste la Descarga Masiva, aquí llegarán automáticamente los documentos.",
    )
    
    # --- NUEVO: listas basadas en rutas + meta ---
    pdf_input_paths: List[str] = []
    docx_input_paths: List[str] = []
    pptx_input_paths: List[str] = []
    unsupported_office: List[Dict[str, str]] = []
    
    docx_meta: Dict[str, Dict[str, Any]] = {}
    pptx_meta: Dict[str, Dict[str, Any]] = {}
    pdf_meta: Dict[str, Dict[str, Any]] = {}
    
    # Directorio base para guardar archivos subidos manualmente
    manual_dir = st.session_state.get("pipeline_manual_dir")
    if manual_dir is None:
        manual_dir = tempfile.mkdtemp(prefix="utp_pipeline_manual_")
        st.session_state["pipeline_manual_dir"] = manual_dir
    manual_dir_path = Path(manual_dir)
    
    # --- 2.1 Documentos provenientes de Descarga Masiva ---
    resultados_desc = st.session_state.get("descarga_resultados") or []
    for r in resultados_desc:
        ruta = r.get("ruta_archivo")
        if not ruta:
            continue
        ruta_str = str(ruta)
        ext = Path(ruta_str).suffix.lower()
        url_or_source = r.get("url")
        
        if ext == ".pdf":
            pdf_input_paths.append(ruta_str)
            pdf_meta[ruta_str] = {
                "display_name": Path(ruta_str).name,
                "source_url": url_or_source,
                "origin": "bulk_download",
            }
        elif ext == ".docx":
            docx_input_paths.append(ruta_str)
            docx_meta[ruta_str] = {
                "display_name": Path(ruta_str).name,
                "source_url": url_or_source,
                "origin": "bulk_download",
            }
        elif ext == ".pptx":
            pptx_input_paths.append(ruta_str)
            pptx_meta[ruta_str] = {
                "display_name": Path(ruta_str).name,
                "source_url": url_or_source,
                "origin": "bulk_download",
            }
        elif ext in (".doc", ".ppt"):
            unsupported_office.append({
                "Archivo": Path(ruta_str).name,
                "Motivo": f"Formato {ext} no soportado. Convierte a .docx o .pptx para poder analizar los links.",
            })
    
    # --- 2.2 Documentos / ZIP subidos manualmente ---
    if uploaded_files:
        for f in uploaded_files:
            fname = f.name
            ext = Path(fname).suffix.lower()
            dest_path = manual_dir_path / Path(fname).name
            
            def _write_if_needed(target: Path, data_bytes: bytes):
                try:
                    with open(target, "wb") as fw:
                        fw.write(data_bytes)
                except Exception as exc:
                    st.warning(f"No se pudo guardar el archivo '{target.name}' en disco: {exc}")
            
            if ext in (".pdf", ".docx", ".pptx", ".doc", ".ppt"):
                data_bytes = f.getbuffer()
                
                if ext == ".pdf":
                    _write_if_needed(dest_path, data_bytes)
                    pdf_input_paths.append(str(dest_path))
                    pdf_meta[str(dest_path)] = {
                        "display_name": Path(fname).name,
                        "source_url": fname,
                        "origin": "upload",
                    }
                elif ext == ".docx":
                    _write_if_needed(dest_path, data_bytes)
                    docx_input_paths.append(str(dest_path))
                    docx_meta[str(dest_path)] = {
                        "display_name": Path(fname).name,
                        "source_url": fname,
                        "origin": "upload",
                    }
                elif ext == ".pptx":
                    _write_if_needed(dest_path, data_bytes)
                    pptx_input_paths.append(str(dest_path))
                    pptx_meta[str(dest_path)] = {
                        "display_name": Path(fname).name,
                        "source_url": fname,
                        "origin": "upload",
                    }
                elif ext in (".doc", ".ppt"):
                    unsupported_office.append({
                        "Archivo": fname,
                        "Motivo": f"Formato {ext} no soportado. Convierte a .docx o .pptx antes de cargarlo.",
                    })
            
            elif ext == ".zip":
                try:
                    zdata = BytesIO(f.getbuffer())
                    with zipfile.ZipFile(zdata, "r") as zf:
                        for info in zf.infolist():
                            if info.is_dir():
                                continue
                            inner_name = Path(info.filename).name
                            inner_ext = Path(inner_name).suffix.lower()
                            file_bytes = zf.read(info)
                            inner_dest = manual_dir_path / inner_name
                            _write_if_needed(inner_dest, file_bytes)
                            
                            if inner_ext == ".pdf":
                                pdf_input_paths.append(str(inner_dest))
                                pdf_meta[str(inner_dest)] = {
                                    "display_name": inner_name,
                                    "source_url": fname,
                                    "origin": "upload_zip",
                                }
                            elif inner_ext == ".docx":
                                docx_input_paths.append(str(inner_dest))
                                docx_meta[str(inner_dest)] = {
                                    "display_name": inner_name,
                                    "source_url": fname,
                                    "origin": "upload_zip",
                                }
                            elif inner_ext == ".pptx":
                                pptx_input_paths.append(str(inner_dest))
                                pptx_meta[str(inner_dest)] = {
                                    "display_name": inner_name,
                                    "source_url": fname,
                                    "origin": "upload_zip",
                                }
                            elif inner_ext in (".doc", ".ppt"):
                                unsupported_office.append({
                                    "Archivo": inner_name,
                                    "Motivo": f"Formato {inner_ext} no soportado. Convierte a .docx o .pptx para poder analizar los links.",
                                })
                except Exception as e:
                    st.warning(f"No se pudo leer el ZIP `{fname}`: {e}")
    
    has_docs = bool(pdf_input_paths or docx_input_paths or pptx_input_paths)
    
    step_pdf1.markdown(
        render_step_header_html(
            "4",
            "Agregar documentos (PDF, Word, PPT) directos o desde ZIP",
            "ok" if has_docs else "warn",
        ),
        unsafe_allow_html=True,
    )
    
    if not has_docs:
        st.caption("Agrega documentos manualmente o ejecuta primero la Descarga Masiva para que lleguen aquí automáticamente.")
        ui_card_close()
        return
    
    # --- Firma de documentos para controlar re-ejecuciones ---
    all_paths_for_signature: List[str] = []
    all_paths_for_signature.extend(pdf_input_paths)
    all_paths_for_signature.extend(docx_input_paths)
    all_paths_for_signature.extend(pptx_input_paths)
    
    signature: List[Tuple[str, int]] = []
    for p in all_paths_for_signature:
        try:
            size = os.path.getsize(p)
        except OSError:
            size = 0
        signature.append((str(Path(p).name), int(size)))
    signature.sort()
    
    prev_sig = st.session_state.get("pipeline_pdf_signature")
    if prev_sig != signature:
        st.session_state["pipeline_pdf_signature"] = signature
        st.session_state["pipeline_pdf_done"] = False
        st.session_state["pipeline_pdf_results"] = None
        st.session_state["pipeline_pdf_errors"] = None
        
        # RESET de rutas generadas y reporte de links / status
        st.session_state["pipeline_docx_paths"] = None
        st.session_state["pipeline_docx_meta"] = None
        st.session_state["pipeline_pptx_paths"] = None
        st.session_state["pipeline_pptx_meta"] = None
        st.session_state["pipeline_word_done"] = False
        st.session_state["pipeline_df_links"] = None
        st.session_state["pipeline_word_errors"] = None
        
        st.session_state["pipeline_status_done"] = False
        st.session_state["status_result_df"] = None
        st.session_state["status_export_df"] = None
        st.session_state["status_invalid_df"] = None
    
    # --- Listado de archivos seleccionados ---
    def _build_pdf_table_from_paths(paths: List[str]) -> pd.DataFrame:
        rows = []
        for p in paths:
            path_obj = Path(p)
            try:
                size_mb = path_obj.stat().st_size / (1024 * 1024)
            except OSError:
                size_mb = 0.0
            pages = "?"
            if fitz is not None and path_obj.exists():
                try:
                    with fitz.open(path_obj) as doc_pdf:
                        pages = len(doc_pdf)
                except Exception:
                    pages = "?"
            rows.append({
                "Nombre": path_obj.name,
                "Tamaño_MB": round(size_mb, 2),
                "Páginas": pages,
            })
        return pd.DataFrame(rows)
    
    def _build_office_table_from_paths(paths: List[str], meta: Dict[str, Dict[str, Any]]) -> pd.DataFrame:
        rows = []
        for p in paths:
            path_obj = Path(p)
            try:
                size_mb = path_obj.stat().st_size / (1024 * 1024)
            except OSError:
                size_mb = 0.0
            m = meta.get(p, {})
            display = m.get("display_name", path_obj.name)
            origen = m.get("origin", "")
            rows.append({
                "Nombre": display,
                "Tamaño_MB": round(size_mb, 2),
                "Origen": origen,
            })
        return pd.DataFrame(rows)
    
    if pdf_input_paths:
        with st.expander("Archivos PDF seleccionados", expanded=False):
            df_files = _build_pdf_table_from_paths(pdf_input_paths)
            st.dataframe(df_files, use_container_width=True, height=260)
    
    if docx_input_paths:
        with st.expander("Archivos Word (DOCX) seleccionados", expanded=False):
            df_files_word = _build_office_table_from_paths(docx_input_paths, docx_meta)
            st.dataframe(df_files_word, use_container_width=True, height=260)
    
    if pptx_input_paths:
        with st.expander("Archivos PPTX seleccionados", expanded=False):
            df_files_ppt = _build_office_table_from_paths(pptx_input_paths, pptx_meta)
            st.dataframe(df_files_ppt, use_container_width=True, height=260)
    
    if unsupported_office:
        with st.expander("⚠️ Archivos Office no soportados", expanded=False):
            df_unsup = pd.DataFrame(unsupported_office)
            st.dataframe(df_unsup, use_container_width=True, height=260)
    
    # Opciones de procesamiento
    with st.expander("Opciones de procesamiento de PDFs", expanded=False):
        col1, col2 = st.columns(2)
        with col1:
            usar_multihilo = st.toggle(
                "Usar procesamiento paralelo por páginas (solo PDFs)",
                value=True,
                help="Recomendado cuando los PDFs tienen muchas páginas.",
                key="extr_usar_multihilo",
            )
        with col2:
            max_workers = st.number_input(
                "Número máximo de workers para PDFs",
                min_value=1,
                max_value=16,
                value=4,
                step=1,
                key="extr_max_workers",
            )
    
    # 4. Procesar todos los documentos (Paso 5)
    render_simple_step_header("5", "Procesar todos los documentos (PDF, Word, PPT)")
    
    progress_pdf_task = st.empty()

    
    auto_trigger_pdf = not st.session_state.get("pipeline_pdf_done", False)
    
    if auto_trigger_pdf:
        try:
            render_task_progress(
                progress_pdf_task,
                "Procesando documentos (PDF → Word)",
                0.0,
                "Iniciando procesamiento de documentos...",
            )

            resultados_pdf: List[Dict[str, Any]] = []
            errores_pdf: List[Dict[str, Any]] = []

            if pdf_input_paths:
                total_pdf_files = len(pdf_input_paths)
                with st.spinner("Extrayendo texto y generando archivos Word desde PDFs..."):
                    for idx, pdf_path in enumerate(pdf_input_paths, start=1):
                        try:
                            processor = PDFBatchProcessor(max_workers=int(max_workers))
                            options = {
                                "usar_multihilo": bool(usar_multihilo),
                                "max_workers": int(max_workers),
                                "filtrar_bibliografia": False,
                            }
                            tmp_dir = tempfile.mkdtemp(prefix="utp_pdf_extr_")
                            result = processor.process_single_pdf(pdf_path, tmp_dir, options)
                            resultados_pdf.append(result)
                        except Exception as e:
                            errores_pdf.append({
                                "status": "error",
                                "file": pdf_path,
                                "error": str(e),
                                "success": False,
                            })

                        pct = idx / total_pdf_files if total_pdf_files > 0 else 1.0
                        detail = f"{idx}/{total_pdf_files} PDFs procesados"
                        render_task_progress(
                            progress_pdf_task,
                            "Procesando documentos (PDF → Word)",
                            pct,
                            detail,
                        )
            else:
                # No hay PDFs; marcamos la tarea como completada
                render_task_progress(
                    progress_pdf_task,
                    "Procesando documentos (PDF → Word)",
                    1.0,
                    "No hay PDFs pendientes; se continuará con Word/PPTX.",
                )


            
            # 2) DOCX generados desde PDFs + DOCX originales (rutas)
            generated_docx_paths: List[str] = []
            generated_docx_meta: Dict[str, Dict[str, Any]] = {}
            
            for r in resultados_pdf or []:
                if r.get("status") == "success":
                    out_path = r.get("output")
                    if out_path and os.path.exists(out_path):
                        out_path_str = str(out_path)
                        original_pdf_name = r.get("file", "")
                        source_url = pdf_meta.get(original_pdf_name, {}).get("source_url")
                        display_stem = Path(original_pdf_name).stem
                        display_name = f"{display_stem}.docx"
                        
                        generated_docx_paths.append(out_path_str)
                        generated_docx_meta[out_path_str] = {
                            "display_name": display_name,
                            "source_url": source_url,
                            "origin": "pdf_to_word",
                        }
            
            combined_docx_paths = list(dict.fromkeys(docx_input_paths + generated_docx_paths))
            combined_docx_meta = dict(docx_meta)
            combined_docx_meta.update(generated_docx_meta)
            
            st.session_state["pipeline_docx_paths"] = combined_docx_paths
            st.session_state["pipeline_docx_meta"] = combined_docx_meta
            st.session_state["pipeline_pptx_paths"] = pptx_input_paths
            st.session_state["pipeline_pptx_meta"] = pptx_meta
            
            st.session_state["pipeline_pdf_results"] = resultados_pdf
            st.session_state["pipeline_pdf_errors"] = errores_pdf
            st.session_state["extraccion_resultados"] = resultados_pdf
            st.session_state["extraccion_errores"] = errores_pdf
            st.session_state["pipeline_pdf_done"] = True
            st.session_state["pipeline_word_inputs_count"] = len(combined_docx_paths)
            st.session_state["pipeline_ppt_inputs_count"] = len(pptx_input_paths)
            
            # progress_bar_pdf.empty()
            # status_text_pdf.empty()
        except Exception as e:
            # progress_bar_pdf.empty()
            # status_text_pdf.empty()
            st.error(f"Ocurrió un error durante el procesamiento de documentos: {e}")
    else:
        resultados_pdf = st.session_state.get("pipeline_pdf_results") or []
        errores_pdf = st.session_state.get("pipeline_pdf_errors") or []
        word_count = st.session_state.get("pipeline_word_inputs_count", 0)
        ppt_count = st.session_state.get("pipeline_ppt_inputs_count", 0)
        
        if resultados_pdf or errores_pdf or word_count or ppt_count:
            total_ok = sum(1 for r in resultados_pdf if r.get("status") == "success")
            total_err = len(errores_pdf)
            total_pdf_files = total_ok + total_err
            
            m1, m2, m3, m4, m5 = st.columns(5)
            m1.metric("PDF procesados", total_pdf_files)
            m2.metric("PDF OK", total_ok)
            m3.metric("PDF con error", total_err)
            m4.metric("Word (DOCX) detectados", word_count)
            m5.metric("PPTX detectados", ppt_count)
    
    if st.session_state.get("pipeline_pdf_done", False):
        render_success_chip("Procesamiento de documentos completado")
    
    ui_card_close()
    
    # ======================================================
    # 3. GrammarScan - Ortografía y Gramática
    # ======================================================
    render_hero(
        "UTP GrammarScan — Ortografía y Gramática (PDF, DOCX, PPTX, TXT)",
        "Herramienta inteligente desarrollada para el análisis y revisión automatizada de ortografía y gramática de diversos tipos de archivos académicos.",
        "📂",
    )
    
    # Parámetros (expander colapsado por defecto)
    with st.expander("Parámetros", expanded=False):
        c1, c2, c3 = st.columns([1, 1, 1])
        with c1:
            lang_code = st.selectbox("Idioma", ["es", "en-US", "pt-BR", "fr", "de"], index=0, key="gs_lang")
        with c2:
            max_chars_call = st.number_input(
                "Máx. caracteres por llamada (LOCAL)",
                3000, 40000, 10000,  # ⬅️ rango y valor por defecto más conservadores
                help="Se agrupan páginas/diapos hasta este límite para mantener contexto.",
                key="gs_max_chars",
            )
        with c3:
            workers = st.slider(
                "Trabajadores (hilos)",
                1, max(2, os.cpu_count() or 4),
                min(2, (os.cpu_count() or 4)),  # ⬅️ por defecto 2
                help="Paraleliza el troceo por páginas. Las llamadas a LT se serializan para estabilidad.",
                key="gs_workers",
            )
    
    with st.expander("Estado del motor", expanded=False):
        st.write(f"Java detectado: **{find_java()}**")
        st.write("Backend LanguageTool: **local** (una sola instancia por sesión).")
        if not find_java():
            st.error("No se puede continuar sin Java.")
            st.stop()
    
    # Paso 6: subir archivos
    ui_card_open()
    step6_ph = st.empty()
    
    # Combinar archivos de diferentes fuentes
    all_uploaded_files = []
    
    # 1. Archivos DOCX generados del pipeline
    docx_paths = st.session_state.get("pipeline_docx_paths") or []
    for docx_path in docx_paths:
        # Crear un objeto similar a UploadedFile
        class FakeUploadedFile:
            def __init__(self, path, display_name):
                self.name = display_name
                self.path = path
                
            def getvalue(self):
                with open(self.path, 'rb') as f:
                    return f.read()
                
            @property
            def size(self):
                return os.path.getsize(self.path)
        
        meta = st.session_state.get("pipeline_docx_meta", {}).get(docx_path, {})
        display_name = meta.get("display_name", Path(docx_path).name)
        fake_file = FakeUploadedFile(docx_path, display_name)
        all_uploaded_files.append(fake_file)
    
    # 2. Archivos PPTX del pipeline
    pptx_paths = st.session_state.get("pipeline_pptx_paths") or []
    for pptx_path in pptx_paths:
        class FakeUploadedFile:
            def __init__(self, path, display_name):
                self.name = display_name
                self.path = path
                
            def getvalue(self):
                with open(self.path, 'rb') as f:
                    return f.read()
                
            @property
            def size(self):
                return os.path.getsize(self.path)
        
        meta = st.session_state.get("pipeline_pptx_meta", {}).get(pptx_path, {})
        display_name = meta.get("display_name", Path(pptx_path).name)
        fake_file = FakeUploadedFile(pptx_path, display_name)
        all_uploaded_files.append(fake_file)
    
    # 3. Archivos PDF originales (opcional, si se quieren analizar directamente)
    # Podríamos permitir analizar PDFs directamente también
    
    # Uploader adicional para archivos manuales
    uploader_key = f"gs_uploader_{st.session_state['gs_uploader_key']}"
    manual_ups = st.file_uploader(
        "6 Sube uno o varios archivos (.pdf, .docx, .pptx, .txt, .zip)",
        type=["pdf", "docx", "pptx", "txt", "zip"],
        accept_multiple_files=True,
        key=uploader_key,
        help="También puedes subir archivos adicionales manualmente.",
    )
    
    if manual_ups:
        all_uploaded_files.extend(manual_ups)
    
    have_files = bool(all_uploaded_files)
    step6_ph.markdown(
        render_step_header_html("6", "Sube uno o varios archivos (.pdf, .docx, .pptx, .txt, .zip)", have_files),
        unsafe_allow_html=True,
    )
    
    if have_files:
        st.success(f"{len(all_uploaded_files)} archivo(s) seleccionado(s).")
        st.caption(
            "Documentos del pipeline automático + archivos subidos manualmente. "
            "Para lotes muy grandes es más estable subir archivos .zip."
        )
    else:
        st.info("Sube al menos un archivo para iniciar el análisis automático.")
    
    ui_card_close()
    
    # Recuperar estado previo
    final_df = st.session_state.get("gs_final_df")
    resumen_completo_df = st.session_state.get("gs_resumen_completo_df")
    metrics = st.session_state.get("gs_metrics")
    elapsed = st.session_state.get("gs_elapsed", 0.0)
    
    if have_files:
        # Firma de archivos
        files_sig = build_files_signature(all_uploaded_files)

        # Firma de parámetros de análisis (para forzar reproceso cuando cambian)
        param_sig = (
            st.session_state.get("gs_lang", "es"),
            st.session_state.get("gs_max_chars", 30000),
            st.session_state.get("gs_workers", 4),
            st.session_state.get("gs_excluir_biblio", True),
            st.session_state.get("gs_modismos", False),
        )

        signature = (files_sig, param_sig)
        last_signature = st.session_state.get("gs_last_files_signature")

        need_processing = (final_df is None) or (signature != last_signature)

        if need_processing:
            final_df, resumen_completo_df, metrics, elapsed = process_grammarscan_files(
                ups=all_uploaded_files,
                lang_code=st.session_state.get("gs_lang", "es"),
                max_chars_call=st.session_state.get("gs_max_chars", 30000),
                workers=st.session_state.get("gs_workers", 4),
                excluir_biblio=st.session_state.get("gs_excluir_biblio", True),
                analizar_modismos=st.session_state.get("gs_modismos", False),
            )
            st.session_state["gs_final_df"] = final_df
            st.session_state["gs_resumen_completo_df"] = resumen_completo_df
            st.session_state["gs_metrics"] = metrics
            st.session_state["gs_elapsed"] = elapsed
            st.session_state["gs_last_files_signature"] = signature

    # Paso 7: Procesar documentos (automático) + resultados
    ui_card_open()
    render_simple_step_header("7", "Procesar documentos (análisis automático)")
    
    if not have_files:
        st.info("Aún no hay análisis. Sube archivos en el Paso 6 para iniciar el proceso automáticamente.")
    else:
        if metrics:
            cA, cB, cC, cD = st.columns(4)
            with cA:
                st.metric("Seleccionados", metrics.get("total", 0))
            with cB:
                st.metric("Con incidencias", metrics.get("n_inc", 0))
            with cC:
                st.metric("Sin incidencias / sin texto", metrics.get("n_zero", 0))
            with cD:
                st.metric("Errores", metrics.get("n_err", 0))
        
        if final_df is not None and not final_df.empty:
            st.subheader("📑 Resultados (detalle de incidencias)")
            st.dataframe(final_df, use_container_width=True, hide_index=True)
        else:
            st.info("No se encontraron incidencias en los archivos procesados.")
        
        if elapsed:
            st.caption(f"⏱️ Tiempo total: {elapsed:0.2f}s")
    
    ui_card_close()

    # Paso 8: Excel (Resultados + Resúmenes)
    ui_card_open()
    render_simple_step_header("8", "Excel (Resultados + Resúmenes)")
    
    if final_df is None or resumen_completo_df is None or resumen_completo_df.empty:
        st.info("Aún no hay resultados para exportar. Procesa documentos primero.")
        # Si no hay datos, limpiamos estado de Excel para futuros corridas
        st.session_state["gs_excel_bytes"] = None
        st.session_state["gs_excel_autotrigger_done"] = False
    else:
        # Generar bytes del Excel y guardarlos en sesión
        excel_bytes = to_excel_bytes(final_df, resumen_completo_df)
        st.session_state["gs_excel_bytes"] = excel_bytes

        # Botón manual (por si el usuario quiere descargar de nuevo)
        st.download_button(
            "⬇️ Excel (Resultados + Resúmenes)",
            data=excel_bytes,
            file_name="UTP_GrammarScan_Resultados.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True,
        )

        # Auto-descarga solo la primera vez que se genera el Excel
        if not st.session_state.get("gs_excel_autotrigger_done", False):
            try:
                b64 = base64.b64encode(excel_bytes).decode("utf-8")
                file_name = "UTP_GrammarScan_Resultados.xlsx"

                # Componente HTML invisible que dispara la descarga al cargarse
                components.html(
                    f"""
                    <html>
                        <body>
                            <a id="auto_download_link"
                               download="{file_name}"
                               href="data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{b64}">
                            </a>
                            <script>
                                document.getElementById('auto_download_link').click();
                            </script>
                        </body>
                    </html>
                    """,
                    height=0,
                    width=0,
                )

                st.session_state["gs_excel_autotrigger_done"] = True
            except Exception as e:
                st.warning(f"No se pudo iniciar la descarga automática del Excel: {e}")
    
    ui_card_close()

# ======================================================
# MAIN
# ======================================================
def main():
    st.set_page_config(
        page_title=APP_TITLE,
        page_icon=APP_ICON,
        layout="wide",
        initial_sidebar_state="expanded",
    )
    
    apply_global_styles()
    init_session_state()
    
    with st.sidebar:
        render_sidebar_header()
        
        st.radio(
            "Módulos",
            MODULES,
            index=MODULES.index(st.session_state["module"]),
            key="module_radio",
            on_change=on_change_module,
        )
        module = st.session_state["module"]
        
        st.markdown("---")
        with st.expander("Recomendaciones"):
            st.markdown(
                """
                - Se recomienda en el proceso "1" descargar de forma masiva entre **500 - 700** registros como máximo.  
                - Para más de **700** registros, lo recomendable es:  
                    Ejecutar el app en local o dividir el Excel de Url en varios archivos (por ejemplo bloques de 500 o 700 registros) y procesarlos por partes.  
                - Esto para evitar que el contenedor de **Streamlit Cloud** se quede sin memoria (~1 GB de RAM) 
                """
            )
    
    module = st.session_state["module"]
    
    if module == "Home":
        page_home()
    elif module == "Report GrammarScan":
        render_report_grammarscan()
    else:
        render_hero(title=module, subtitle="Módulo no encontrado.", icon="⚠️")
        st.error("Módulo seleccionado no existe.")

if __name__ == "__main__":
    main()






























